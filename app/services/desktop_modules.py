from __future__ import annotations

from datetime import date, datetime, time, timedelta
import hashlib
import html
import json
import mimetypes
from pathlib import Path
import re
import shutil
import sqlite3
import uuid
import urllib.error
import urllib.parse
import urllib.request

from sqlalchemy import func, or_, text

from app import db
from app.models import (
    AIChangeSet, AIConversation, AIMessage, ActivityEvent, ApiPreset, CalendarEvent,
    Executor, FileOperation, LabRecord, LabRecordRevision, LabRecordStep, LibraryItem, LiteratureItem,
    Note, ResearchProject, SearchDocument, Task, User, WeeklyReport, WeeklyReportEntry,
    WeeklyReportUpdate,
    WorkspaceSetting, ZoteroCollection, ZoteroConnection, lab_record_library_item, lab_record_literature,
    literature_library_item, note_lab_record, note_literature, project_library_item,
    project_literature, project_zotero_collection, utcnow, zotero_collection_literature,
    weekly_report_library_item,
)
from .errors import ConflictError, NotFoundError, ValidationError


LITERATURE_SOURCES = {"manual", "zotero", "import"}
READ_STATUSES = {"unread", "reading", "read"}
NOTE_KINDS = {"idea", "experiment_guide", "literature", "meeting", "supervisor_feedback", "general"}
TASK_STATUSES = {"todo", "doing", "blocked", "done", "cancelled"}
TASK_PRIORITIES = {"low", "medium", "high"}
FILE_KINDS = {"raw_data", "image", "spreadsheet", "protocol", "report", "literature", "other"}
STORAGE_MODES = {"managed", "external", "zotero"}
WEEKLY_FILE_EXTENSIONS = {
    ".csv", ".doc", ".docx", ".key", ".md", ".odp", ".ods", ".pdf",
    ".ppt", ".pptx", ".txt", ".xls", ".xlsx", ".zip",
}
AI_TARGET_FIELDS = {
    "project": {"title", "code", "objective", "status", "start_date", "end_date", "notes"},
    "lab_record": {
        "title", "status", "experiment_date", "executor_snapshot", "location", "objective", "background",
        "hypothesis", "design", "materials_conditions", "expected_result", "actual_process_summary",
        "actual_result", "analysis", "conclusion", "next_steps", "steps",
    },
    "task": {"title", "category", "priority", "deadline", "status", "notes", "project_id", "lab_record_id"},
    "calendar_event": {
        "title", "event_type", "starts_at", "ends_at", "all_day", "notes", "project_id", "lab_record_id",
    },
    "note": {"title", "body", "kind", "project_id"},
    "weekly_report": {
        "title", "report_date", "project_id", "status", "summary", "body", "issues_and_feedback",
        "next_week_plan", "annotation",
    },
}

AI_FIELD_LABELS = {
    "title": "标题", "code": "编号", "objective": "研究目标", "status": "状态", "start_date": "开始日期",
    "end_date": "结束日期", "notes": "说明", "experiment_date": "实验日期", "executor_snapshot": "执行者",
    "location": "地点", "background": "研究背景", "hypothesis": "实验假设", "design": "实验设计",
    "materials_conditions": "材料与条件", "expected_result": "预期结果", "actual_process_summary": "实际过程",
    "actual_result": "实际结果", "analysis": "分析", "conclusion": "结论", "next_steps": "下一步",
    "category": "分类", "priority": "优先级", "deadline": "截止日期", "project_id": "所属项目",
    "lab_record_id": "关联实验记录", "event_type": "事件类型", "starts_at": "开始时间", "ends_at": "结束时间",
    "all_day": "全天事件", "kind": "笔记类型", "report_date": "报告日期", "summary": "摘要",
    "body": "正文", "issues_and_feedback": "问题与反馈", "next_week_plan": "下周计划", "annotation": "周报批注",
    "steps": "实验步骤",
}

AI_DATE_FIELDS = {"start_date", "end_date", "experiment_date", "deadline", "report_date"}
AI_DATETIME_FIELDS = {"starts_at", "ends_at"}
AI_ID_FIELDS = {"project_id", "lab_record_id"}
AI_BOOLEAN_FIELDS = {"all_day"}
AI_ENUM_FIELDS = {
    ("project", "status"): {"active", "paused", "completed", "archived"},
    ("lab_record", "status"): {"draft", "in_progress", "awaiting_analysis", "completed", "archived"},
    ("task", "status"): TASK_STATUSES,
    ("task", "priority"): TASK_PRIORITIES,
    ("note", "kind"): NOTE_KINDS,
    ("calendar_event", "event_type"): {"meeting", "reminder", "custom"},
    ("weekly_report", "status"): {"draft", "submitted", "reviewed", "archived"},
}

AI_FIELD_MAX_LENGTH = {
    "title": 500, "code": 80, "status": 30, "executor_snapshot": 120, "location": 180,
    "category": 40, "priority": 10, "event_type": 30, "kind": 30,
}

AI_TARGET_INSTRUCTIONS = {
    "project": (
        "整理当前研究项目的名称、编号、研究目标、状态、起止日期和说明。"
        "研究目标应具体且可执行；没有明确日期或编号时不要补造。"
    ),
    "lab_record": (
        "当前对象是一条完整实验记录。用户要求生成实验报告时，按研究背景、实验目的、实验假设、"
        "实验设计、材料与条件、实验步骤、预期结果、实际过程、实际结果、分析、结论和下一步分别写入对应字段。"
        "只能根据 current_target、source_snapshot 和用户请求整理实际过程、实际结果、分析与结论；"
        "缺少证据的实际内容必须省略，不能把预期结果改写成实际结果。"
        "steps 是完整的有序步骤数组；修改步骤时必须保留已有步骤的 is_done 和 actual_deviation，"
        "除非用户明确要求改变它们。每项只能包含 title、instruction、planned_duration_minutes、"
        "checkpoint、risk、is_done、actual_deviation。"
    ),
    "task": (
        "把请求整理为一个可执行待办，包括标题、分类、优先级、截止日期、状态、备注和已有对象关联。"
        "截止日期使用 YYYY-MM-DD；关联编号只能沿用上下文中已有的正整数。"
    ),
    "calendar_event": (
        "把安排整理为一个日历事件，包括标题、类型、开始时间、结束时间、全天状态、备注和已有对象关联。"
        "时间使用 ISO 8601；没有明确结束时间时省略 ends_at，不能擅自推测持续时长。"
    ),
    "note": (
        "把内容整理为当前项目下的笔记标题、正文和笔记类型。"
        "保留原始事实与引用边界，不把建议写成已经完成的实验事实。"
    ),
    "weekly_report": (
        "整理当前周报的标题、日期、项目、状态、摘要、正文、问题与反馈和下周计划。"
        "annotation 表示新增一条周报批注、说明或指导，应用时会追加到时间线，不会覆盖已有批注。"
        "只汇总上下文中已有的工作，不得把计划写成已完成事项。"
    ),
}


def build_desktop_ai_system_prompt(target_type, field_schema, *, today, timezone):
    """Build the strict, record-centric prompt shared by every desktop AI provider."""
    if target_type not in AI_TARGET_INSTRUCTIONS:
        raise ValidationError("AI 目标类型无效。")
    allowed = sorted(AI_TARGET_FIELDS[target_type])
    return "\n".join((
        "你是 R/LAB 独立桌面研究助手，负责把用户请求转换成当前对象可直接写入的字段建议。",
        "只输出合法 JSON，不要输出 Markdown 或代码围栏。固定格式："
        '{"reply":"简短说明","proposal":{"字段":"新值"}}。',
        "proposal 只包含确实需要修改且有依据的字段；不要输出空占位，不要重复未变化字段。"
        "除非用户明确要求清空，否则不要用空字符串覆盖已有内容。",
        "不得编造用户未提供的测量数据、日期、人员、编号或完成状态。"
        "用户请求与当前数据冲突时，在 reply 中简要说明，并省略没有可靠依据的字段。",
        "所有字段值必须符合字段 Schema；枚举值必须原样使用，日期和时间必须使用 Schema 指定格式。"
        "project_id 与 lab_record_id 只能使用上下文中已经出现的有效编号。",
        AI_TARGET_INSTRUCTIONS[target_type],
        f"当前日期：{today.isoformat()}。工作区时区：{timezone}。",
        f"目标类型：{target_type}。允许字段：{_json(allowed)}。字段 Schema：{_json(field_schema)}。",
    ))


def _string(value, maximum, field, *, required=False):
    result = str(value or "").strip()
    if required and not result:
        raise ValidationError("请填写必填字段。", field_errors={field: "此字段不能为空"})
    if len(result) > maximum:
        raise ValidationError("字段内容过长。", field_errors={field: f"最多 {maximum} 个字符"})
    return result


def _page_result(query, payload, serializer):
    """Opt-in server pagination while preserving the legacy list response."""
    if not payload.get("pagination"):
        return [serializer(item) for item in query.all()]
    try:
        page = max(1, int(payload.get("page") or 1))
        page_size = min(200, max(1, int(payload.get("page_size") or 50)))
    except (TypeError, ValueError) as exc:
        raise ValidationError("分页参数无效。") from exc
    total = query.order_by(None).count()
    pages = max(1, (total + page_size - 1) // page_size)
    page = min(page, pages)
    rows = query.offset((page - 1) * page_size).limit(page_size).all()
    return {
        "items": [serializer(item) for item in rows],
        "pagination": {"page": page, "page_size": page_size, "pages": pages, "total": total},
    }


def _sorted(query, payload, allowed, default):
    key = str(payload.get("sort") or "").strip() or default
    if key not in allowed:
        raise ValidationError("排序参数无效。", field_errors={"sort": "不支持的排序方式"})
    return query.order_by(*allowed[key])


def _positive_id(value, field="id", *, optional=False):
    if optional and value in (None, ""):
        return None
    try:
        result = int(value)
    except (TypeError, ValueError) as exc:
        raise ValidationError("编号无效。", field_errors={field: "必须为整数"}) from exc
    if result <= 0:
        raise ValidationError("编号无效。", field_errors={field: "必须大于 0"})
    return result


def _date(value, field, *, default=None):
    if value in (None, ""):
        return default
    try:
        return date.fromisoformat(str(value)[:10])
    except ValueError as exc:
        raise ValidationError("日期格式无效。", field_errors={field: "请使用 YYYY-MM-DD"}) from exc


def _datetime(value, field):
    if not value:
        raise ValidationError("时间不能为空。", field_errors={field: "此字段不能为空"})
    try:
        return datetime.fromisoformat(str(value).replace("Z", "+00:00")).replace(tzinfo=None)
    except ValueError as exc:
        raise ValidationError("时间格式无效。", field_errors={field: "请使用 ISO 8601 时间"}) from exc


def _iso(value):
    return value.isoformat() if value else None


def _json(value):
    return json.dumps(value, ensure_ascii=False, separators=(",", ":"))


def _load_json(value, fallback):
    try:
        parsed = json.loads(value or "")
        return parsed if isinstance(parsed, type(fallback)) else fallback
    except (TypeError, json.JSONDecodeError):
        return fallback


def _normalize_doi(value):
    doi = str(value or "").strip()
    return re.sub(r"^(?:https?://(?:dx\.)?doi\.org/|doi:\s*)", "", doi, flags=re.I).lower()


def _zotero_select_url(source_key):
    """Build a Zotero select link without leaking library parsing into the UI."""
    parts = str(source_key or "").split(":")
    if len(parts) == 3 and parts[0] == "zotero" and parts[1].isdigit() and parts[2]:
        return f"zotero://select/library/items/{parts[2]}"
    if (
        len(parts) == 4
        and parts[0] == "zotero"
        and parts[1] in {"group", "groups"}
        and parts[2].isdigit()
        and parts[3]
    ):
        return f"zotero://select/groups/{parts[2]}/items/{parts[3]}"
    return ""


def _zotero_response_header(response, name):
    headers = getattr(response, "headers", None)
    if headers is not None and hasattr(headers, "get"):
        return headers.get(name)
    getheader = getattr(response, "getheader", None)
    return getheader(name) if callable(getheader) else None


def _fetch_zotero_top_items(base_url, *, page_size=100):
    """Fetch every top-level item before source-missing reconciliation."""
    items = []
    start = 0
    total_results = None
    previous_page_signature = None
    while True:
        request = urllib.request.Request(
            f"{base_url}/api/users/0/items/top?limit={page_size}&start={start}&format=json",
            headers={"Accept": "application/json", "Zotero-API-Version": "3"},
            method="GET",
        )
        with urllib.request.urlopen(request, timeout=8) as response:
            page = json.loads(response.read().decode("utf-8"))
            total_header = _zotero_response_header(response, "Total-Results")
        if not isinstance(page, list):
            raise ValueError("response is not a list")
        page_signature = tuple(
            str(raw.get("key") or raw.get("data", {}).get("key") or "")
            if isinstance(raw, dict) and isinstance(raw.get("data", {}), dict)
            else ""
            for raw in page
        )
        if start and page and page_signature == previous_page_signature:
            raise ValueError("Zotero pagination returned the same page twice")
        previous_page_signature = page_signature
        if total_header not in (None, ""):
            try:
                total_results = max(0, int(total_header))
            except (TypeError, ValueError):
                total_results = None
        items.extend(page)
        if not page or len(page) < page_size:
            break
        if total_results is not None and len(items) >= total_results:
            break
        start += len(page)
        if start > 1_000_000:
            raise ValueError("Zotero pagination exceeded the safety limit")
    return items


class ZoteroSyncCancelled(Exception):
    pass


def _fetch_zotero_items(base_url, library_path, *, since=None, page_size=100, cancel_event=None, resource="items"):
    """Fetch all changed items and return response identity/version metadata."""
    items = []
    start = 0
    total_results = None
    previous_page_signature = None
    server_id = ""
    library_version = 0
    while True:
        if cancel_event is not None and cancel_event.is_set():
            raise ZoteroSyncCancelled()
        query = {"limit": page_size, "start": start, "format": "json"}
        if resource == "items":
            query["includeTrashed"] = 1
        if since is not None:
            query["since"] = int(since)
        request = urllib.request.Request(
            f"{base_url}/api/{library_path}/{resource}?{urllib.parse.urlencode(query)}",
            headers={"Accept": "application/json", "Zotero-API-Version": "3"},
            method="GET",
        )
        with urllib.request.urlopen(request, timeout=8) as response:
            page = json.loads(response.read().decode("utf-8"))
            total_header = _zotero_response_header(response, "Total-Results")
            server_id = _zotero_response_header(response, "Zotero-Server-ID") or server_id
            version_header = _zotero_response_header(response, "Last-Modified-Version")
        if not isinstance(page, list):
            raise ValueError("response is not a list")
        page_signature = tuple(
            str(raw.get("key") or raw.get("data", {}).get("key") or "")
            if isinstance(raw, dict) and isinstance(raw.get("data", {}), dict)
            else ""
            for raw in page
        )
        if start and page and page_signature == previous_page_signature:
            raise ValueError("Zotero pagination returned the same page twice")
        previous_page_signature = page_signature
        if total_header not in (None, ""):
            try:
                total_results = max(0, int(total_header))
            except (TypeError, ValueError):
                total_results = None
        if version_header not in (None, ""):
            try:
                library_version = max(library_version, int(version_header))
            except (TypeError, ValueError):
                pass
        items.extend(page)
        if not page or len(page) < page_size:
            break
        if total_results is not None and len(items) >= total_results:
            break
        start += len(page)
        if start > 1_000_000:
            raise ValueError("Zotero pagination exceeded the safety limit")
    return items, {
        "server_id": str(server_id or ""),
        "library_version": library_version,
        "complete_library": since is None,
    }


def _sha256(path):
    digest = hashlib.sha256()
    with Path(path).open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


RECORD_EXPORT_STATUS_LABELS = {
    "draft": "草稿",
    "in_progress": "进行中",
    "awaiting_analysis": "待分析",
    "completed": "已完成",
    "archived": "已归档",
}
RECORD_EXPORT_PHASE_LABELS = {
    "planned": "计划值",
    "actual": "实际值",
    "derived": "推导值",
}
RECORD_EXPORT_EVENT_LABELS = {
    "note": "观察记录",
    "measurement": "测量记录",
    "status": "状态变更",
    "checkpoint": "检查点",
}


def _export_value(value, fallback="未填写"):
    """Return a printable value without inventing experimental facts."""
    if value is None:
        return fallback
    if isinstance(value, bool):
        return "是" if value else "否"
    text_value = str(value).strip()
    return text_value or fallback


def _export_date(value):
    return _iso(value) or "未设置"


def _record_export_view(record, attachment_rows=()):
    """Build the format-neutral content model used by Word, PDF and Markdown.

    The desktop editor is the source of truth for this model.  Optional legacy
    relations (events, parameters, materials and files) are included only when
    they actually exist, so an empty record remains an honest, useful template.
    """
    project = record.project
    status_label = RECORD_EXPORT_STATUS_LABELS.get(record.status, _export_value(record.status))
    conclusion = _export_value(record.conclusion, "")
    result_label = "已记录" if conclusion else "待确认"
    metadata = [
        ("实验编号", _export_value(record.record_code, "未设置")),
        ("记录编号", f"R-{record.id}"),
        ("实验日期", _export_date(record.experiment_date)),
        ("实验人员", _export_value(record.executor_snapshot)),
        ("所属项目", _export_value(project.title)),
        ("项目编号", _export_value(project.code, "未设置")),
        ("地点或设备", _export_value(record.location)),
        ("结果判定", result_label),
        ("记录状态", status_label),
    ]
    sections = [
        ("02", "背景、目的与假设", (
            ("实验目的", record.objective),
            ("研究背景与依据", record.background),
            ("实验假设", record.hypothesis),
        )),
        ("03", "实验设计与材料条件", (
            ("实验设计", record.design),
            ("材料与条件", record.materials_conditions),
        )),
        ("05", "预期结果", (
            ("预期观察、数据范围与成功标准", record.expected_result),
        )),
        ("06", "实际过程与结果", (
            ("实际过程摘要", record.actual_process_summary),
            ("实际结果与数据摘要", record.actual_result),
        )),
        ("07", "分析、结论与下一步", (
            ("分析", record.analysis),
            ("结论", record.conclusion),
            ("下一步", record.next_steps),
        )),
    ]
    steps = []
    for index, step in enumerate(record.steps, 1):
        duration = f"{step.planned_duration_minutes} 分钟" if step.planned_duration_minutes is not None else "未设置"
        steps.append({
            "index": index,
            "title": _export_value(step.title),
            "instruction": _export_value(step.instruction),
            "duration": duration,
            "checkpoint": _export_value(step.checkpoint),
            "risk": _export_value(step.risk),
            "status": "已完成" if step.is_done else "未完成",
            "deviation": _export_value(step.actual_deviation),
            "executor": _export_value(step.executor_snapshot),
        })
    events = [
        {
            "date": _iso(event.occurred_at) or "未设置",
            "type": RECORD_EXPORT_EVENT_LABELS.get(event.event_type, _export_value(event.event_type)),
            "executor": _export_value(event.executor_snapshot),
            "content": _export_value(event.content),
        }
        for event in record.events
    ]
    parameters = [
        {
            "phase": RECORD_EXPORT_PHASE_LABELS.get(parameter.phase, _export_value(parameter.phase)),
            "name": _export_value(parameter.name),
            "value": _export_value(parameter.value_text),
            "unit": _export_value(parameter.unit),
            "notes": _export_value(parameter.notes),
        }
        for parameter in record.parameters
    ]
    materials = [
        {
            "kind": _export_value(material.kind),
            "name": _export_value(material.name),
            "identifier": _export_value(material.identifier),
            "role": _export_value(material.role),
            "planned": _export_value(material.planned_amount),
            "actual": _export_value(material.actual_amount),
            "unit": _export_value(material.unit),
        }
        for material in record.materials
    ]
    attachments = [
        {
            "name": _export_value(row.display_name),
            "kind": _export_value(row.kind),
            "mode": _export_value(row.storage_mode),
            "status": _export_value(row.link_status),
        }
        for row in attachment_rows
    ]
    return {
        "title": _export_value(record.title),
        "record_code": _export_value(record.record_code, "未设置"),
        "project_title": _export_value(project.title),
        "project_code": _export_value(project.code, "未设置"),
        "experiment_date": _export_date(record.experiment_date),
        "metadata": metadata,
        "sections": sections,
        "steps": steps,
        "events": events,
        "parameters": parameters,
        "materials": materials,
        "attachments": attachments,
        "status_label": status_label,
    }


def _export_markup(value):
    """Escape a record value for a ReportLab Paragraph while keeping line breaks."""
    escaped = html.escape(_export_value(value)).replace("\r\n", "\n").replace("\r", "\n")
    return escaped.replace("\n", "<br/>")


def _export_section_groups(view):
    """Keep the first-page workflow and continuation-page results together."""
    before = tuple(section for section in view["sections"] if section[0] in {"02", "03"})
    after = tuple(section for section in view["sections"] if section[0] not in {"02", "03"})
    return before, after


def _write_record_markdown(target, view):
    before_steps, after_steps = _export_section_groups(view)
    lines = [
        "# R/LAB · 实验记录报告", "", f"# {view['title']}", "",
        f"**记录编号**：{view['metadata'][1][1]}  ",
        f"**实验编号**：{view['record_code']}  ",
        f"**所属项目**：{view['project_title']}（{view['project_code']}）  ",
        f"**实验日期**：{view['experiment_date']}  ",
        f"**记录状态**：{view['status_label']}", "", "## 01 基本信息", "",
    ]
    lines.extend(f"- **{label}**：{value}" for label, value in view["metadata"])
    for number, heading, fields in before_steps:
        lines.extend(["", f"## {number} {heading}", ""])
        for label, value in fields:
            lines.extend([f"### {label}", "", _export_value(value), ""])
    lines.extend(["## 04 实验步骤", ""])
    if view["steps"]:
        for step in view["steps"]:
            lines.extend([
                f"### {step['index']}. {step['title']}（{step['status']}）", "", step["instruction"],
                f"\n- 时长：{step['duration']}\n- 检查点：{step['checkpoint']}\n- 风险：{step['risk']}\n- 执行者：{step['executor']}\n- 实际偏差：{step['deviation']}", "",
            ])
    else:
        lines.append("未填写")
    lines.extend(["", "---"])
    for number, heading, fields in after_steps:
        lines.extend(["", f"## {number} {heading}", ""])
        for label, value in fields:
            lines.extend([f"### {label}", "", _export_value(value), ""])
    if view["events"]:
        lines.extend(["## 过程与观察记录", "", "| 日期 | 类型 | 执行者 | 内容 |", "| --- | --- | --- | --- |"])
        lines.extend(
            f"| {item['date']} | {item['type']} | {item['executor']} | {item['content'].replace('|', '\\|')} |"
            for item in view["events"]
        )
    lines.extend(["", "## 08 结构化参数", ""])
    if view["parameters"]:
        lines.extend(["| 阶段 | 参数 | 数值 | 单位 | 说明 |", "| --- | --- | --- | --- | --- |"])
        lines.extend(
            f"| {item['phase']} | {item['name']} | {item['value']} | {item['unit']} | {item['notes'].replace('|', '\\|')} |"
            for item in view["parameters"]
        )
    else:
        lines.append("当前记录暂无结构化参数。")
    if view["materials"]:
        lines.extend(["", "## 09 结构化材料", "", "| 类型 | 名称 | 标识 | 用途 | 计划量 | 实际量 | 单位 |", "| --- | --- | --- | --- | --- | --- | --- |"])
        lines.extend(
            f"| {item['kind']} | {item['name']} | {item['identifier']} | {item['role']} | {item['planned']} | {item['actual']} | {item['unit']} |"
            for item in view["materials"]
        )
    lines.extend(["", "## 10 原始数据与附件", ""])
    if view["attachments"]:
        lines.extend(f"- {item['name']}（{item['kind']}，{item['mode']}，{item['status']}）" for item in view["attachments"])
    else:
        lines.append("当前记录没有关联文件；如需查看文件，请返回软件中的“文件与数据”页。")
    lines.extend(["", "> 人工核验提示：导出内容仅整理当前记录中已保存的字段，不会把预期结果改写成实际结果。"])
    target.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _write_record_docx(target, view):
    from docx import Document
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm, Pt, RGBColor

    document = Document()
    section = document.sections[0]
    section.top_margin, section.bottom_margin = Cm(1.8), Cm(1.7)
    section.left_margin, section.right_margin = Cm(1.9), Cm(1.9)
    accent, ink = "5B5CE2", "232430"

    def set_font(run, size=None, *, bold=None, color=None, name="思源宋体"):
        run.font.name = name
        run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), name)
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.bold = bold
        if color:
            run.font.color.rgb = RGBColor.from_string(color)

    def style_font(style, size, *, bold=False, color=ink):
        style.font.name = "思源宋体"
        style._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), "思源宋体")
        style.font.size = Pt(size)
        style.font.bold = bold
        style.font.color.rgb = RGBColor.from_string(color)

    for style_name, size, bold, color in (
        ("Normal", 10.5, False, ink), ("Title", 22, True, ink),
        ("Heading 1", 13, True, ink), ("Heading 2", 11, True, ink),
    ):
        style = document.styles[style_name]
        style_font(style, size, bold=bold, color=color)
        style.paragraph_format.space_after = Pt(6)

    def shade_cell(cell, fill):
        tc_pr = cell._tc.get_or_add_tcPr()
        shading = tc_pr.find(qn("w:shd"))
        if shading is None:
            shading = OxmlElement("w:shd")
            tc_pr.append(shading)
        shading.set(qn("w:fill"), fill)

    def table_borders(table, color="D7D8E2", size="6"):
        tbl_pr = table._tbl.tblPr
        borders = tbl_pr.first_child_found_in("w:tblBorders")
        if borders is None:
            borders = OxmlElement("w:tblBorders")
            tbl_pr.append(borders)
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            tag = "w:" + edge
            element = borders.find(qn(tag))
            if element is None:
                element = OxmlElement(tag)
                borders.append(element)
            element.set(qn("w:val"), "single")
            element.set(qn("w:sz"), size)
            element.set(qn("w:color"), color)

    def repeat_header(row):
        tr_pr = row._tr.get_or_add_trPr()
        header = OxmlElement("w:tblHeader")
        header.set(qn("w:val"), "true")
        tr_pr.append(header)

    def put_cell(cell, value, *, bold=False, color=ink, fill=None, size=9.2):
        cell.text = ""
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        if fill:
            shade_cell(cell, fill)
        paragraph = cell.paragraphs[0]
        paragraph.paragraph_format.space_after = Pt(0)
        run = paragraph.add_run(_export_value(value))
        set_font(run, size, bold=bold, color=color)

    def add_section_heading(number, heading):
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.space_before = Pt(9)
        paragraph.paragraph_format.space_after = Pt(5)
        paragraph.paragraph_format.keep_with_next = True
        run = paragraph.add_run(f"{number}  ")
        set_font(run, 10, bold=True, color=accent)
        run = paragraph.add_run(heading)
        set_font(run, 13, bold=True, color=ink)
        p_pr = paragraph._p.get_or_add_pPr()
        borders = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "5")
        bottom.set(qn("w:space"), "5")
        bottom.set(qn("w:color"), "D7D8E2")
        borders.append(bottom)
        p_pr.append(borders)

    def add_field(label, value):
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(4)
        paragraph.paragraph_format.keep_together = True
        run = paragraph.add_run(label)
        set_font(run, 10, bold=True, color=accent)
        run = paragraph.add_run("\n" + _export_value(value))
        set_font(run, 10.5, color=ink)

    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    set_font(header.add_run("R/LAB · 实验记录报告"), 8.5, color="6B6D7A")
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_font(footer.add_run("R/LAB  ·  第 "), 8, color="7C7E8B")
    field_run = footer.add_run()
    set_font(field_run, 8, color="7C7E8B")
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    field_run._r.extend((begin, instr, end))
    set_font(footer.add_run(" 页"), 8, color="7C7E8B")

    document.core_properties.title = view["title"]
    document.core_properties.subject = "R/LAB 实验记录报告"
    title = document.add_paragraph(style="Title")
    title_run = title.add_run(view["title"])
    set_font(title_run, 22, bold=True, color=ink)
    subtitle = document.add_paragraph()
    subtitle.paragraph_format.space_after = Pt(10)
    set_font(subtitle.add_run(
        f"{view['record_code']}  ·  {view['project_title']}（{view['project_code']}）  ·  {view['experiment_date']}"
    ), 10, color="6B6D7A")

    add_section_heading("01", "基本信息")
    metadata_table = document.add_table(rows=0, cols=4)
    metadata_table.alignment = WD_TABLE_ALIGNMENT.LEFT
    metadata_table.autofit = False
    for offset in range(0, len(view["metadata"]), 2):
        row = metadata_table.add_row()
        pairs = view["metadata"][offset:offset + 2]
        while len(pairs) < 2:
            pairs.append(("", ""))
        for pair_index, (label, value) in enumerate(pairs):
            label_cell = row.cells[pair_index * 2]
            value_cell = row.cells[pair_index * 2 + 1]
            label_cell.width, value_cell.width = Cm(2.2), Cm(6.0)
            put_cell(label_cell, label, bold=True, color=accent, fill="F1F1FD", size=9)
            put_cell(value_cell, value, size=9.2)
    table_borders(metadata_table)

    before_steps, after_steps = _export_section_groups(view)
    for number, heading, fields in before_steps:
        add_section_heading(number, heading)
        for label, value in fields:
            add_field(label, value)

    add_section_heading("04", "实验步骤")
    if view["steps"]:
        steps_table = document.add_table(rows=1, cols=6)
        steps_table.alignment = WD_TABLE_ALIGNMENT.LEFT
        steps_table.autofit = False
        headers = ("序号", "步骤 / 阶段", "操作说明", "时长", "状态", "检查点 / 风险 / 偏差")
        for cell, text_value in zip(steps_table.rows[0].cells, headers):
            put_cell(cell, text_value, bold=True, color="FFFFFF", fill=accent, size=8.5)
        repeat_header(steps_table.rows[0])
        for step in view["steps"]:
            row = steps_table.add_row()
            values = (
                step["index"], step["title"], step["instruction"], step["duration"], step["status"],
                f"检查点：{step['checkpoint']}\n风险：{step['risk']}\n偏差：{step['deviation']}",
            )
            for cell, value in zip(row.cells, values):
                put_cell(cell, value, size=8.5)
        table_borders(steps_table)
    else:
        add_field("实验步骤", "未填写")

    document.add_page_break()
    for number, heading, fields in after_steps:
        add_section_heading(number, heading)
        for label, value in fields:
            add_field(label, value)

    if view["events"]:
        add_section_heading("06A", "过程与观察记录")
        events_table = document.add_table(rows=1, cols=4)
        events_table.alignment = WD_TABLE_ALIGNMENT.LEFT
        events_table.autofit = False
        for cell, text_value in zip(events_table.rows[0].cells, ("日期", "类型", "执行者", "内容")):
            put_cell(cell, text_value, bold=True, color="FFFFFF", fill=accent, size=8.5)
        repeat_header(events_table.rows[0])
        for item in view["events"]:
            row = events_table.add_row()
            for cell, value in zip(row.cells, (item["date"], item["type"], item["executor"], item["content"])):
                put_cell(cell, value, size=8.5)
        table_borders(events_table)

    add_section_heading("08", "结构化参数")
    if view["parameters"]:
        parameter_table = document.add_table(rows=1, cols=5)
        parameter_table.alignment = WD_TABLE_ALIGNMENT.LEFT
        parameter_table.autofit = False
        for cell, text_value in zip(parameter_table.rows[0].cells, ("阶段", "参数", "数值", "单位", "说明")):
            put_cell(cell, text_value, bold=True, color="FFFFFF", fill=accent, size=8.5)
        repeat_header(parameter_table.rows[0])
        for item in view["parameters"]:
            row = parameter_table.add_row()
            for cell, value in zip(row.cells, (item["phase"], item["name"], item["value"], item["unit"], item["notes"])):
                put_cell(cell, value, size=8.5)
        table_borders(parameter_table)
    else:
        add_field("参数记录", "当前记录暂无结构化参数。")

    if view["materials"]:
        add_section_heading("09", "结构化材料")
        material_table = document.add_table(rows=1, cols=7)
        material_table.alignment = WD_TABLE_ALIGNMENT.LEFT
        material_table.autofit = False
        for cell, text_value in zip(material_table.rows[0].cells, ("类型", "名称", "标识", "用途", "计划量", "实际量", "单位")):
            put_cell(cell, text_value, bold=True, color="FFFFFF", fill=accent, size=8.2)
        repeat_header(material_table.rows[0])
        for item in view["materials"]:
            row = material_table.add_row()
            for cell, value in zip(row.cells, (item["kind"], item["name"], item["identifier"], item["role"], item["planned"], item["actual"], item["unit"])):
                put_cell(cell, value, size=8.2)
        table_borders(material_table)

    add_section_heading("10", "原始数据与附件")
    if view["attachments"]:
        for item in view["attachments"]:
            add_field(item["name"], f"{item['kind']} · {item['mode']} · {item['status']}")
    else:
        add_field("附件", "当前记录没有关联文件；如需查看文件，请返回软件中的“文件与数据”页。")
    note_table = document.add_table(rows=1, cols=1)
    note_table.alignment = WD_TABLE_ALIGNMENT.LEFT
    put_cell(
        note_table.cell(0, 0),
        "人工核验提示：导出内容仅整理当前记录中已保存的字段，不会把预期结果改写成实际结果。",
        color="5D6070", fill="F0F1F5", size=9,
    )
    table_borders(note_table, color="E4E6EC", size="3")
    document.save(target)


def _write_record_pdf(target, view, font_path):
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    font_name = "RLabCJK"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        try:
            pdfmetrics.registerFont(TTFont(font_name, font_path))
        except Exception as exc:
            raise ValidationError("中文字体无法加载，PDF 导出不可用。") from exc
    ink_color = colors.HexColor("#232430")
    accent_color = colors.HexColor("#5B5CE2")
    muted_color = colors.HexColor("#6B6D7A")
    styles = {
        "kicker": ParagraphStyle("rlab-kicker", fontName=font_name, fontSize=8.5, leading=11, textColor=muted_color, spaceAfter=3),
        "title": ParagraphStyle("rlab-title", fontName=font_name, fontSize=19, leading=25, textColor=ink_color, spaceAfter=4),
        "subtitle": ParagraphStyle("rlab-subtitle", fontName=font_name, fontSize=9.5, leading=14, textColor=muted_color, spaceAfter=10),
        "section": ParagraphStyle("rlab-section", fontName=font_name, fontSize=12, leading=17, textColor=ink_color, spaceBefore=9, spaceAfter=6, keepWithNext=True),
        "label": ParagraphStyle("rlab-label", fontName=font_name, fontSize=9, leading=13, textColor=accent_color, spaceAfter=2),
        "body": ParagraphStyle("rlab-body", fontName=font_name, fontSize=9.5, leading=15, textColor=ink_color, spaceAfter=6, wordWrap="CJK"),
        "table": ParagraphStyle("rlab-table", fontName=font_name, fontSize=8.2, leading=11.5, textColor=ink_color, wordWrap="CJK"),
        "table_head": ParagraphStyle("rlab-table-head", fontName=font_name, fontSize=8.2, leading=11, textColor=colors.white, wordWrap="CJK"),
        "note": ParagraphStyle("rlab-note", fontName=font_name, fontSize=8.5, leading=13, textColor=muted_color, spaceBefore=8),
    }

    def para(value, style="body"):
        return Paragraph(_export_markup(value), styles[style])

    def heading(number, text_value):
        return Paragraph(f'<font color="#5B5CE2">{html.escape(number)}</font>  {html.escape(text_value)}', styles["section"])

    def fields(label, value):
        return [Paragraph(html.escape(label), styles["label"]), para(value)]

    def styled_table(rows, widths):
        table = Table(rows, colWidths=widths, repeatRows=1, hAlign="LEFT")
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), accent_color),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), .35, colors.HexColor("#D7D8E2")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FAFAFD")]),
        ]))
        return table

    metadata_rows = []
    metadata_without_record = [view["metadata"][0], *view["metadata"][2:]]
    for offset in range(0, len(metadata_without_record), 2):
        pairs = list(metadata_without_record[offset:offset + 2])
        while len(pairs) < 2:
            pairs.append(("", ""))
        metadata_rows.append([
            Paragraph(html.escape(pairs[0][0]), styles["label"]), para(pairs[0][1], "table"),
            Paragraph(html.escape(pairs[1][0]), styles["label"]), para(pairs[1][1], "table"),
        ])
    metadata_table = Table(metadata_rows, colWidths=[24 * mm, 54 * mm, 24 * mm, 54 * mm], hAlign="LEFT")
    metadata_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#F1F1FD")),
        ("BACKGROUND", (2, 0), (2, -1), colors.HexColor("#F1F1FD")),
        ("GRID", (0, 0), (-1, -1), .35, colors.HexColor("#D7D8E2")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))

    story = [
        Paragraph("R/LAB · 实验记录报告", styles["kicker"]),
        Paragraph(html.escape(view["title"]), styles["title"]),
        Paragraph(html.escape(f"{view['record_code']}  ·  {view['project_title']}（{view['project_code']}）  ·  {view['experiment_date']}"), styles["subtitle"]),
        heading("01", "基本信息"), metadata_table,
    ]
    before_steps, after_steps = _export_section_groups(view)
    for number, heading_text, section_fields in before_steps:
        story.append(heading(number, heading_text))
        for label, value in section_fields:
            story.extend(fields(label, value))
    story.append(heading("04", "实验步骤"))
    if view["steps"]:
        rows = [[Paragraph(html.escape(value), styles["table_head"]) for value in ("序号", "步骤 / 阶段", "操作说明", "时长", "状态", "检查点 / 风险 / 偏差")]]
        for step in view["steps"]:
            combined = f"检查点：{step['checkpoint']}<br/>风险：{step['risk']}<br/>偏差：{step['deviation']}"
            rows.append([para(value, "table") for value in (step["index"], step["title"], step["instruction"], step["duration"], step["status"], combined)])
        story.append(styled_table(rows, [10 * mm, 25 * mm, 38 * mm, 16 * mm, 16 * mm, 63 * mm]))
    else:
        story.extend(fields("实验步骤", "未填写"))
    story.append(PageBreak())
    for number, heading_text, section_fields in after_steps:
        story.append(heading(number, heading_text))
        for label, value in section_fields:
            story.extend(fields(label, value))
    if view["events"]:
        story.append(heading("06A", "过程与观察记录"))
        rows = [[Paragraph(html.escape(value), styles["table_head"]) for value in ("日期", "类型", "执行者", "内容")]]
        for item in view["events"]:
            rows.append([para(item[key], "table") for key in ("date", "type", "executor", "content")])
        story.append(styled_table(rows, [34 * mm, 25 * mm, 28 * mm, 81 * mm]))
    story.append(heading("08", "结构化参数"))
    if view["parameters"]:
        rows = [[Paragraph(html.escape(value), styles["table_head"]) for value in ("阶段", "参数", "数值", "单位", "说明")]]
        for item in view["parameters"]:
            rows.append([para(item[key], "table") for key in ("phase", "name", "value", "unit", "notes")])
        story.append(styled_table(rows, [22 * mm, 31 * mm, 35 * mm, 20 * mm, 60 * mm]))
    else:
        story.extend(fields("参数记录", "当前记录暂无结构化参数。"))
    if view["materials"]:
        story.append(heading("09", "结构化材料"))
        keys = ("kind", "name", "identifier", "role", "planned", "actual", "unit")
        rows = [[Paragraph(html.escape(value), styles["table_head"]) for value in ("类型", "名称", "标识", "用途", "计划量", "实际量", "单位")]]
        for item in view["materials"]:
            rows.append([para(item[key], "table") for key in keys])
        story.append(styled_table(rows, [18 * mm, 31 * mm, 25 * mm, 31 * mm, 22 * mm, 22 * mm, 18 * mm]))
    story.append(heading("10", "原始数据与附件"))
    if view["attachments"]:
        for item in view["attachments"]:
            story.extend(fields(item["name"], f"{item['kind']} · {item['mode']} · {item['status']}"))
    else:
        story.extend(fields("附件", "当前记录没有关联文件；如需查看文件，请返回软件中的“文件与数据”页。"))
    note_table = Table(
        [[Paragraph("人工核验提示：导出内容仅整理当前记录中已保存的字段，不会把预期结果改写成实际结果。", styles["note"])]],
        colWidths=[172 * mm],
        hAlign="LEFT",
    )
    note_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F0F1F5")),
        ("BOX", (0, 0), (-1, -1), .45, colors.HexColor("#E4E6EC")),
        ("LEFTPADDING", (0, 0), (-1, -1), 8), ("RIGHTPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING", (0, 0), (-1, -1), 6), ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(note_table)

    def draw_header_footer(canvas, doc):
        canvas.saveState()
        canvas.setFont(font_name, 8)
        canvas.setFillColor(muted_color)
        canvas.drawString(19 * mm, 10 * mm, "R/LAB · 实验记录报告")
        canvas.drawRightString(191 * mm, 10 * mm, f"第 {doc.page} 页")
        canvas.restoreState()

    document = SimpleDocTemplate(
        str(target), pagesize=A4, leftMargin=19 * mm, rightMargin=19 * mm,
        topMargin=17 * mm, bottomMargin=17 * mm, title=view["title"],
    )
    document.build(story, onFirstPage=draw_header_footer, onLaterPages=draw_header_footer)


def _record_export_named_sections(view):
    """Return the seven editor sections with the report's Chinese numbering."""
    numbers = {
        "02": "二、背景、目的与假设",
        "03": "三、实验设计与材料条件",
        "05": "五、预期结果",
        "06": "六、实际过程与结果",
        "07": "七、分析、结论与下一步",
    }
    return [(numbers.get(number, heading), fields) for number, heading, fields in view["sections"]]


def _write_record_docx_table_template(target, view):
    """Write the clean, table-first Word template used by desktop exports."""
    from docx import Document
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm, Pt, RGBColor

    document = Document()
    section = document.sections[0]
    section.top_margin, section.bottom_margin = Cm(1.35), Cm(1.25)
    section.left_margin, section.right_margin = Cm(1.45), Cm(1.45)
    ink, navy, grid, fill = "202A36", "244A6B", "AEB8C3", "EEF2F5"

    def set_font(run, size=9.5, *, bold=False, color=ink):
        run.font.name = "思源宋体"
        run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), "思源宋体")
        run.font.size = Pt(size)
        run.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)

    def set_cell_shading(cell, color):
        tc_pr = cell._tc.get_or_add_tcPr()
        shading = tc_pr.find(qn("w:shd"))
        if shading is None:
            shading = OxmlElement("w:shd")
            tc_pr.append(shading)
        shading.set(qn("w:fill"), color)

    def set_cell_margins(cell, top=70, start=90, bottom=70, end=90):
        tc = cell._tc
        tc_pr = tc.get_or_add_tcPr()
        margins = tc_pr.first_child_found_in("w:tcMar")
        if margins is None:
            margins = OxmlElement("w:tcMar")
            tc_pr.append(margins)
        for side, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
            node = margins.find(qn(f"w:{side}"))
            if node is None:
                node = OxmlElement(f"w:{side}")
                margins.append(node)
            node.set(qn("w:w"), str(value))
            node.set(qn("w:type"), "dxa")

    def table_borders(table, color=grid, size="5"):
        tbl_pr = table._tbl.tblPr
        borders = tbl_pr.first_child_found_in("w:tblBorders")
        if borders is None:
            borders = OxmlElement("w:tblBorders")
            tbl_pr.append(borders)
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            node = borders.find(qn(f"w:{edge}"))
            if node is None:
                node = OxmlElement(f"w:{edge}")
                borders.append(node)
            node.set(qn("w:val"), "single")
            node.set(qn("w:sz"), size)
            node.set(qn("w:color"), color)

    def repeat_header(row):
        tr_pr = row._tr.get_or_add_trPr()
        header = OxmlElement("w:tblHeader")
        header.set(qn("w:val"), "true")
        tr_pr.append(header)

    def cell_text(cell, value, *, label=False, header=False, size=8.8):
        cell.text = ""
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        set_cell_margins(cell)
        if label:
            set_cell_shading(cell, fill)
        if header:
            set_cell_shading(cell, navy)
        paragraph = cell.paragraphs[0]
        paragraph.paragraph_format.space_after = Pt(0)
        paragraph.paragraph_format.line_spacing = 1.05
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if label or header else WD_ALIGN_PARAGRAPH.LEFT
        run = paragraph.add_run(_export_value(value))
        set_font(run, size=size, bold=label or header, color="FFFFFF" if header else ink)

    def add_heading(text_value):
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.space_before = Pt(7)
        paragraph.paragraph_format.space_after = Pt(4)
        paragraph.paragraph_format.keep_with_next = True
        run = paragraph.add_run(text_value)
        set_font(run, size=11.5, bold=True, color=navy)
        p_pr = paragraph._p.get_or_add_pPr()
        borders = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "4")
        bottom.set(qn("w:color"), navy)
        borders.append(bottom)
        p_pr.append(borders)

    def add_two_col_table(rows):
        table = document.add_table(rows=0, cols=2)
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        table.autofit = False
        for label, value in rows:
            row = table.add_row()
            row.cells[0].width, row.cells[1].width = Cm(4.0), Cm(13.6)
            cell_text(row.cells[0], label, label=True, size=8.8)
            cell_text(row.cells[1], value, size=9.1)
        table_borders(table)
        return table

    def add_table(headers, rows, widths, *, header_size=8.1, body_size=8.2):
        table = document.add_table(rows=1, cols=len(headers))
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        table.autofit = False
        for cell, label, width in zip(table.rows[0].cells, headers, widths):
            cell.width = Cm(width)
            cell_text(cell, label, header=True, size=header_size)
        repeat_header(table.rows[0])
        for values in rows:
            row = table.add_row()
            for cell, value, width in zip(row.cells, values, widths):
                cell.width = Cm(width)
                cell_text(cell, value, size=body_size)
        table_borders(table)
        return table

    # Plain document title: no product name or branding.
    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.paragraph_format.space_after = Pt(3)
    set_font(title.add_run("实验记录报告"), size=18, bold=True, color=ink)
    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.paragraph_format.space_after = Pt(8)
    set_font(subtitle.add_run(view["title"]), size=12, bold=True, color=navy)

    add_heading("一、基本信息")
    metadata = [("实验标题", view["title"])] + list(view["metadata"])
    metadata_table = document.add_table(rows=0, cols=4)
    metadata_table.alignment = WD_TABLE_ALIGNMENT.LEFT
    metadata_table.autofit = False
    first = metadata_table.add_row()
    first.cells[0].width, first.cells[1].width = Cm(2.1), Cm(6.5)
    first.cells[2].width, first.cells[3].width = Cm(2.1), Cm(6.9)
    cell_text(first.cells[0], "实验标题", label=True)
    cell_text(first.cells[1], view["title"], size=9.1)
    cell_text(first.cells[2], "记录编号", label=True)
    cell_text(first.cells[3], view["metadata"][1][1], size=9.1)
    metadata_without_record = [view["metadata"][0], *view["metadata"][2:]]
    for offset in range(0, len(metadata_without_record), 2):
        pairs = list(metadata_without_record[offset:offset + 2])
        if len(pairs) < 2:
            pairs.append(("", ""))
        row = metadata_table.add_row()
        for pair_index, (label, value) in enumerate(pairs):
            cell_text(row.cells[pair_index * 2], label, label=True)
            cell_text(row.cells[pair_index * 2 + 1], value, size=9.1)
    table_borders(metadata_table)

    named_sections = _record_export_named_sections(view)
    before_steps = named_sections[:2]
    after_steps = named_sections[2:]
    for heading, rows in before_steps:
        add_heading(heading)
        add_two_col_table(rows)

    add_heading("四、实验步骤")
    step_rows = []
    for step in view["steps"]:
        step_rows.append((
            step["index"], step["title"], step["instruction"], step["duration"],
            f"{step['status']}\n{step['executor']}", step["checkpoint"],
            f"风险：{step['risk']}\n偏差：{step['deviation']}",
        ))
    if not step_rows:
        step_rows = [("", "未填写", "当前记录暂无实验步骤。", "", "", "", "")]
    add_table(
        ("序号", "步骤 / 阶段", "操作说明", "计划时长", "状态 / 人员", "检查点", "风险 / 实际偏差"),
        step_rows,
        (0.9, 2.25, 4.8, 1.55, 1.8, 2.5, 4.2),
    )

    document.add_page_break()
    for heading, rows in after_steps:
        add_heading(heading)
        add_two_col_table(rows)

    if view["events"]:
        add_heading("过程与观察记录")
        add_table(
            ("日期", "类型", "执行者", "记录内容"),
            [(item["date"], item["type"], item["executor"], item["content"]) for item in view["events"]],
            (3.1, 2.5, 3.0, 8.8),
        )

    add_heading("附、结构化参数")
    parameter_rows = [
        (item["name"], item["value"], "", item["unit"], item["notes"])
        for item in view["parameters"]
    ] or [("当前记录暂无参数", "", "", "", "")]
    add_table(("参数名称", "计划值", "实际值", "单位", "备注"), parameter_rows, (3.5, 3.1, 3.1, 2.0, 4.8))

    add_heading("附、原始数据与附件")
    attachment_rows = [
        (item["name"], item["kind"], item["mode"], item["status"], "")
        for item in view["attachments"]
    ] or [("当前记录暂无关联文件", "", "", "", "请在软件的文件与数据页查看")]
    add_table(("文件名称", "类型", "来源", "状态", "备注"), attachment_rows, (5.0, 2.3, 2.7, 2.5, 4.0))

    add_heading("人工核验提示")
    review = document.add_table(rows=1, cols=1)
    review.alignment = WD_TABLE_ALIGNMENT.LEFT
    cell_text(
        review.cell(0, 0),
        "请核对实验编号、日期、人员、关键参数、结果数据及附件是否完整、准确、可追溯；如与原始记录不符，请在软件中修订。",
        size=8.8,
    )
    set_cell_shading(review.cell(0, 0), "F5F6F7")
    table_borders(review, color="C8D0D8", size="4")

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_font(footer.add_run("第 "), size=8, color="6A7480")
    field_run = footer.add_run()
    set_font(field_run, size=8, color="6A7480")
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    field_run._r.extend((begin, instr, end))
    set_font(footer.add_run(" 页"), size=8, color="6A7480")
    document.core_properties.title = view["title"]
    document.core_properties.subject = "实验记录"
    document.save(target)


def _write_record_pdf_table_template(target, view, font_path):
    """Write the clean table-first PDF counterpart of the Word template."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Table, TableStyle

    font_name = "RecordCJK"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        try:
            pdfmetrics.registerFont(TTFont(font_name, font_path))
        except Exception as exc:
            raise ValidationError("中文字体无法加载，PDF 导出不可用。") from exc
    ink = colors.HexColor("#202A36")
    navy = colors.HexColor("#244A6B")
    line = colors.HexColor("#AEB8C3")
    light = colors.HexColor("#EEF2F5")
    styles = {
        "title": ParagraphStyle("record-title", fontName=font_name, fontSize=17, leading=22, alignment=1, textColor=ink, spaceAfter=2),
        "subtitle": ParagraphStyle("record-subtitle", fontName=font_name, fontSize=11, leading=15, alignment=1, textColor=navy, spaceAfter=8),
        "section": ParagraphStyle("record-section", fontName=font_name, fontSize=10.8, leading=15, textColor=navy, spaceBefore=6, spaceAfter=4, keepWithNext=True),
        "label": ParagraphStyle("record-label", fontName=font_name, fontSize=8.0, leading=10.5, textColor=ink, alignment=1, wordWrap="CJK"),
        "body": ParagraphStyle("record-body", fontName=font_name, fontSize=8.3, leading=11.2, textColor=ink, wordWrap="CJK"),
        "head": ParagraphStyle("record-head", fontName=font_name, fontSize=7.8, leading=10, textColor=colors.white, alignment=1, wordWrap="CJK"),
        "note": ParagraphStyle("record-note", fontName=font_name, fontSize=8.1, leading=11.5, textColor=ink, wordWrap="CJK"),
    }

    def para(value, style="body"):
        return Paragraph(_export_markup(value), styles[style])

    def heading(text_value):
        return Paragraph(html.escape(text_value), styles["section"])

    def table(rows, widths, *, repeat_rows=0, label_columns=()):
        result = Table(rows, colWidths=[width * mm for width in widths], repeatRows=repeat_rows, hAlign="LEFT")
        commands = [
            ("GRID", (0, 0), (-1, -1), .35, line),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4), ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]
        for column in label_columns:
            commands.append(("BACKGROUND", (column, 0), (column, -1), light))
        if repeat_rows:
            commands.extend([
                ("BACKGROUND", (0, 0), (-1, repeat_rows - 1), navy),
                ("TEXTCOLOR", (0, 0), (-1, repeat_rows - 1), colors.white),
            ])
        result.setStyle(TableStyle(commands))
        return result

    def two_col(rows):
        return table(
            [[Paragraph(html.escape(label), styles["label"]), para(value)] for label, value in rows],
            (39, 133),
            label_columns=(0,),
        )

    def grid(headers, rows, widths):
        all_rows = [[Paragraph(html.escape(value), styles["head"]) for value in headers]]
        all_rows.extend([para(value) for value in row] for row in rows)
        return table(all_rows, widths, repeat_rows=1)

    story = [Paragraph("实验记录报告", styles["title"]), Paragraph(html.escape(view["title"]), styles["subtitle"])]
    story.append(heading("一、基本信息"))
    metadata_rows = [[
        Paragraph("实验标题", styles["label"]), para(view["title"]),
        Paragraph("记录编号", styles["label"]), para(view["metadata"][1][1]),
    ]]
    metadata_without_record = [view["metadata"][0], *view["metadata"][2:]]
    for offset in range(0, len(metadata_without_record), 2):
        pairs = list(metadata_without_record[offset:offset + 2])
        if len(pairs) < 2:
            pairs.append(("", ""))
        metadata_rows.append([
            Paragraph(html.escape(pairs[0][0]), styles["label"]), para(pairs[0][1], "body"),
            Paragraph(html.escape(pairs[1][0]), styles["label"]), para(pairs[1][1], "body"),
        ])
    story.append(table(metadata_rows, (24, 47, 24, 77), label_columns=(0, 2)))

    named_sections = _record_export_named_sections(view)
    for section_heading, rows in named_sections[:2]:
        story.extend([heading(section_heading), two_col(rows)])

    story.append(heading("四、实验步骤"))
    step_rows = [(
        str(step["index"]), step["title"], step["instruction"], step["duration"],
        f"{step['status']}\n{step['executor']}", step["checkpoint"],
        f"风险：{step['risk']}\n偏差：{step['deviation']}",
    ) for step in view["steps"]]
    if not step_rows:
        step_rows = [("", "未填写", "当前记录暂无实验步骤。", "", "", "", "")]
    story.append(grid(
        ("序号", "步骤 / 阶段", "操作说明", "计划时长", "状态 / 人员", "检查点", "风险 / 实际偏差"),
        step_rows,
        (8, 20, 39, 15, 18, 24, 56),
    ))
    story.append(PageBreak())

    for section_heading, rows in named_sections[2:]:
        story.extend([heading(section_heading), two_col(rows)])
    if view["events"]:
        story.extend([
            heading("过程与观察记录"),
            grid(
                ("日期", "类型", "执行者", "记录内容"),
                [(item["date"], item["type"], item["executor"], item["content"]) for item in view["events"]],
                (33, 26, 30, 89),
            ),
        ])

    story.append(heading("附、结构化参数"))
    parameter_rows = [
        (item["name"], "", item["value"], item["unit"], item["notes"])
        for item in view["parameters"]
    ] or [("当前记录暂无参数", "", "", "", "")]
    story.append(grid(("参数名称", "计划值", "实际值", "单位", "备注"), parameter_rows, (35, 31, 31, 21, 60)))

    story.append(heading("附、原始数据与附件"))
    attachment_rows = [
        (item["name"], item["kind"], item["mode"], item["status"], "")
        for item in view["attachments"]
    ] or [("当前记录暂无关联文件", "", "", "", "请在软件的文件与数据页查看")]
    story.append(grid(("文件名称", "类型", "来源", "状态", "备注"), attachment_rows, (50, 24, 28, 26, 50)))
    story.append(heading("人工核验提示"))
    review = Table([[Paragraph("请核对实验编号、日期、人员、关键参数、结果数据及附件是否完整、准确、可追溯；如与原始记录不符，请在软件中修订。", styles["note"])]], colWidths=[174 * mm], hAlign="LEFT")
    review.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F5F6F7")),
        ("BOX", (0, 0), (-1, -1), .4, colors.HexColor("#C8D0D8")),
        ("LEFTPADDING", (0, 0), (-1, -1), 7), ("RIGHTPADDING", (0, 0), (-1, -1), 7),
        ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(review)

    def draw_page_number(canvas, doc):
        canvas.saveState()
        canvas.setFont(font_name, 8)
        canvas.setFillColor(colors.HexColor("#6A7480"))
        canvas.drawCentredString(105 * mm, 8 * mm, f"第 {doc.page} 页")
        canvas.restoreState()

    document = SimpleDocTemplate(
        str(target), pagesize=A4, leftMargin=15 * mm, rightMargin=15 * mm,
        topMargin=12 * mm, bottomMargin=13 * mm, title=view["title"],
    )
    document.build(story, onFirstPage=draw_page_number, onLaterPages=draw_page_number)


def _write_record_markdown_table_template(target, view):
    """Keep Markdown export aligned with the same seven-section table schema."""
    def md(value):
        return _export_value(value).replace("|", "\\|").replace("\n", "<br>")

    lines = ["# 实验记录报告", "", f"## {md(view['title'])}", "", "| 字段 | 内容 |", "| --- | --- |"]
    lines.extend(f"| {md(label)} | {md(value)} |" for label, value in [("实验标题", view["title"]), *view["metadata"]])
    for heading, rows in _record_export_named_sections(view)[:2]:
        lines.extend(["", f"## {heading}", "", "| 字段 | 内容 |", "| --- | --- |"])
        lines.extend(f"| {md(label)} | {md(value)} |" for label, value in rows)
    lines.extend(["", "## 四、实验步骤", "", "| 序号 | 步骤 / 阶段 | 操作说明 | 计划时长 | 状态 / 人员 | 检查点 | 风险 / 实际偏差 |", "| --- | --- | --- | --- | --- | --- | --- |"])
    if view["steps"]:
        lines.extend(
            f"| {md(step['index'])} | {md(step['title'])} | {md(step['instruction'])} | {md(step['duration'])} | {md(step['status'])} / {md(step['executor'])} | {md(step['checkpoint'])} | 风险：{md(step['risk'])}<br>偏差：{md(step['deviation'])} |"
            for step in view["steps"]
        )
    else:
        lines.append("|  | 未填写 | 当前记录暂无实验步骤。 |  |  |  |  |")
    for heading, rows in _record_export_named_sections(view)[2:]:
        lines.extend(["", f"## {heading}", "", "| 字段 | 内容 |", "| --- | --- |"])
        lines.extend(f"| {md(label)} | {md(value)} |" for label, value in rows)
    if view["events"]:
        lines.extend(["", "## 过程与观察记录", "", "| 日期 | 类型 | 执行者 | 记录内容 |", "| --- | --- | --- | --- |"])
        lines.extend(f"| {md(item['date'])} | {md(item['type'])} | {md(item['executor'])} | {md(item['content'])} |" for item in view["events"])
    lines.extend(["", "## 附、结构化参数", "", "| 参数名称 | 计划值 | 实际值 | 单位 | 备注 |", "| --- | --- | --- | --- | --- |"])
    if view["parameters"]:
        lines.extend(f"| {md(item['name'])} |  | {md(item['value'])} | {md(item['unit'])} | {md(item['notes'])} |" for item in view["parameters"])
    else:
        lines.append("| 当前记录暂无参数 |  |  |  |  |")
    lines.extend(["", "## 附、原始数据与附件", "", "| 文件名称 | 类型 | 来源 | 状态 | 备注 |", "| --- | --- | --- | --- | --- |"])
    if view["attachments"]:
        lines.extend(f"| {md(item['name'])} | {md(item['kind'])} | {md(item['mode'])} | {md(item['status'])} |  |" for item in view["attachments"])
    else:
        lines.append("| 当前记录暂无关联文件 |  |  |  | 请在软件的文件与数据页查看 |")
    lines.extend(["", "> 人工核验：请核对实验编号、日期、人员、关键参数、结果数据及附件是否完整、准确、可追溯。"])
    target.write_text("\n".join(lines) + "\n", encoding="utf-8")


class DesktopModuleServiceMixin:
    """Record-centric desktop modules mixed into DesktopApplicationService."""

    @staticmethod
    def _ensure_project(workspace_id, project_id):
        if project_id is None:
            return None
        project = ResearchProject.query.filter_by(id=project_id, workspace_id=workspace_id, is_deleted=False).first()
        if not project:
            raise NotFoundError("所属项目不存在或已移入回收站。")
        return project

    @staticmethod
    def _activity(workspace, event_type, summary, source_table, source_id, *, project_id=None, record_id=None):
        db.session.add(ActivityEvent(
            workspace_id=workspace.id, project_id=project_id, record_id=record_id,
            event_type=event_type, summary=summary[:500], source_table=source_table, source_id=source_id,
        ))

    @staticmethod
    def _index_entity(workspace, entity_type, item):
        if entity_type == "project":
            config = (item.title, f"{item.objective or ''}\n{item.notes or ''}", item.code or "", "projects")
        elif entity_type == "record":
            config = (item.title, "\n".join(str(getattr(item, field) or "") for field in (
                "objective", "background", "hypothesis", "design", "materials_conditions", "expected_result",
                "actual_process_summary", "actual_result", "analysis", "conclusion", "next_steps",
            )), item.record_code, "record-edit")
        elif entity_type == "literature":
            config = (item.title, f"{item.authors_json}\n{item.abstract}\n{item.reading_notes}", f"{item.doi}\n{item.keywords_json}", "literature")
        elif entity_type == "file":
            config = (item.display_name, f"{item.description}\n{item.external_path or item.managed_relative_path or ''}", item.sha256, "files")
        elif entity_type == "note":
            config = (item.title, item.body, item.kind, "notes")
        elif entity_type == "task":
            config = (item.title, item.notes or "", f"{item.category}\n{item.status}", "tasks")
        elif entity_type == "weekly":
            config = (item.title, f"{item.body}\n{item.issues_and_feedback}\n{item.next_week_plan}", item.status, "weekly")
        else:
            return
        document = SearchDocument.query.filter_by(
            workspace_id=workspace.id, entity_type=entity_type, entity_id=item.id,
        ).first()
        if getattr(item, "is_deleted", False):
            if document:
                db.session.delete(document)
            return
        if not document:
            document = SearchDocument(workspace_id=workspace.id, entity_type=entity_type, entity_id=item.id)
            db.session.add(document)
        document.title, document.body, document.keywords, document.view_key = config

    @staticmethod
    def _literature_dto(item, *, detail=False):
        data = {
            "id": item.id, "source": item.source, "source_key": item.source_key,
            "zotero_select_url": _zotero_select_url(item.source_key) if item.source == "zotero" else "",
            "title": item.title, "authors": _load_json(item.authors_json, []), "year": item.year,
            "journal": item.journal, "doi": item.doi, "read_status": item.read_status,
            "source_missing": item.source_missing, "updated_at": _iso(item.updated_at),
            "row_version": item.row_version,
        }
        if detail:
            data["collection_ids"] = [row[0] for row in db.session.execute(
                db.select(zotero_collection_literature.c.collection_id).where(
                    zotero_collection_literature.c.literature_id == item.id
                )
            )]
            project_rows = db.session.execute(
                db.select(ResearchProject.id, ResearchProject.title)
                .select_from(project_literature)
                .join(ResearchProject, ResearchProject.id == project_literature.c.project_id)
                .where(
                    project_literature.c.literature_id == item.id,
                    ResearchProject.is_deleted.is_(False),
                )
                .order_by(ResearchProject.title)
            ).all()
            record_rows = db.session.execute(
                db.select(
                    LabRecord.id,
                    LabRecord.title,
                    LabRecord.project_id,
                    ResearchProject.title,
                    LabRecord.experiment_date,
                )
                .select_from(lab_record_literature)
                .join(LabRecord, LabRecord.id == lab_record_literature.c.record_id)
                .join(ResearchProject, ResearchProject.id == LabRecord.project_id)
                .where(
                    lab_record_literature.c.literature_id == item.id,
                    LabRecord.is_deleted.is_(False),
                    ResearchProject.is_deleted.is_(False),
                )
                .order_by(LabRecord.experiment_date.desc(), LabRecord.updated_at.desc())
            ).all()
            note_rows = db.session.execute(
                db.select(Note.id, Note.title, Note.project_id)
                .select_from(note_literature)
                .join(Note, Note.id == note_literature.c.note_id)
                .where(
                    note_literature.c.literature_id == item.id,
                    Note.is_deleted.is_(False),
                )
                .order_by(Note.updated_at.desc())
            ).all()
            attachment_rows = db.session.execute(
                db.select(
                    LibraryItem.id,
                    LibraryItem.display_name,
                    LibraryItem.storage_mode,
                    LibraryItem.mime_type,
                    LibraryItem.link_status,
                    LibraryItem.source_key,
                )
                .select_from(literature_library_item)
                .join(LibraryItem, LibraryItem.id == literature_library_item.c.library_item_id)
                .where(
                    literature_library_item.c.literature_id == item.id,
                    LibraryItem.is_deleted.is_(False),
                )
                .order_by(LibraryItem.updated_at.desc())
            ).all()
            data.update({
                "volume": item.volume, "issue": item.issue, "pages": item.pages, "url": item.url,
                "abstract": item.abstract, "keywords": _load_json(item.keywords_json, []),
                "reading_notes": item.reading_notes, "synced_at": _iso(item.synced_at),
                "projects": [
                    {"id": project_id, "title": project_title}
                    for project_id, project_title in project_rows
                ],
                "records": [
                    {
                        "id": record_id,
                        "title": record_title,
                        "project_id": project_id,
                        "project_title": project_title,
                        "experiment_date": _iso(experiment_date),
                    }
                    for record_id, record_title, project_id, project_title, experiment_date in record_rows
                ],
                "notes": [
                    {"id": note_id, "title": note_title, "project_id": project_id}
                    for note_id, note_title, project_id in note_rows
                ],
                "attachments": [
                    {
                        "id": attachment_id,
                        "display_name": display_name,
                        "storage_mode": storage_mode,
                        "mime_type": mime_type,
                        "link_status": link_status,
                        "source_key": source_key,
                        "zotero_select_url": _zotero_select_url(source_key)
                        if storage_mode == "zotero" else "",
                    }
                    for attachment_id, display_name, storage_mode, mime_type, link_status, source_key
                    in attachment_rows
                ],
                "project_ids": [row[0] for row in project_rows],
                "record_ids": [row[0] for row in record_rows],
            })
        return data

    def list_literature(self, payload=None):
        payload = payload or {}
        search = _string(payload.get("search"), 200, "search")
        source = _string(payload.get("source"), 20, "source")
        read_status = _string(payload.get("read_status"), 20, "read_status")
        project_filter = payload.get("project_id")
        collection_filter = _positive_id(payload.get("collection_id"), "collection_id", optional=True)

        def query(workspace):
            items = LiteratureItem.query.filter_by(workspace_id=workspace.id, is_deleted=False)
            if source:
                items = items.filter_by(source=source)
            if read_status:
                items = items.filter_by(read_status=read_status)
            if project_filter == "unclassified":
                linked_ids = db.select(project_literature.c.literature_id)
                items = items.filter(~LiteratureItem.id.in_(linked_ids))
            elif project_filter not in (None, ""):
                project_id = _positive_id(project_filter, "project_id")
                self._ensure_project(workspace.id, project_id)
                direct_ids = db.select(project_literature.c.literature_id).where(project_literature.c.project_id == project_id)
                mapped_collections = db.select(project_zotero_collection.c.collection_id).where(
                    project_zotero_collection.c.project_id == project_id
                )
                mapped_ids = db.select(zotero_collection_literature.c.literature_id).where(
                    zotero_collection_literature.c.collection_id.in_(mapped_collections)
                )
                items = items.filter(or_(LiteratureItem.id.in_(direct_ids), LiteratureItem.id.in_(mapped_ids)))
            if collection_filter:
                collection = ZoteroCollection.query.filter_by(id=collection_filter, workspace_id=workspace.id).first()
                if not collection: raise NotFoundError("Zotero Collection 不存在。")
                member_ids = db.select(zotero_collection_literature.c.literature_id).where(
                    zotero_collection_literature.c.collection_id == collection_filter
                )
                items = items.filter(LiteratureItem.id.in_(member_ids))
            if search:
                pattern = f"%{search.replace('%', r'\%').replace('_', r'\_')}%"
                items = items.filter(or_(
                    LiteratureItem.title.ilike(pattern, escape="\\"),
                    LiteratureItem.authors_json.ilike(pattern, escape="\\"),
                    LiteratureItem.journal.ilike(pattern, escape="\\"),
                    LiteratureItem.doi.ilike(pattern, escape="\\"),
                ))
            items = _sorted(items, payload, {
                "updated_desc": (LiteratureItem.updated_at.desc(),),
                "title_asc": (LiteratureItem.title.asc(),),
                "year_desc": (LiteratureItem.year.desc(), LiteratureItem.title.asc()),
            }, "updated_desc")
            if payload.get("pagination"):
                try:
                    page = max(1, int(payload.get("page") or 1)); page_size = min(200, max(1, int(payload.get("page_size") or 50)))
                except (TypeError, ValueError) as exc: raise ValidationError("分页参数无效。") from exc
                total = items.order_by(None).count(); pages = max(1, (total + page_size - 1) // page_size); page = min(page, pages)
                rows = items.offset((page - 1) * page_size).limit(page_size).all()
            else:
                rows = items.all(); total = len(rows); page = 1; page_size = max(1, total); pages = 1
            ids = [item.id for item in rows]
            project_counts = dict(db.session.execute(
                db.select(project_literature.c.literature_id, func.count()).where(
                    project_literature.c.literature_id.in_(ids)
                ).group_by(project_literature.c.literature_id)
            ).all()) if ids else {}
            record_counts = dict(db.session.execute(
                db.select(lab_record_literature.c.literature_id, func.count()).where(
                    lab_record_literature.c.literature_id.in_(ids)
                ).group_by(lab_record_literature.c.literature_id)
            ).all()) if ids else {}
            result = []
            for item in rows:
                data = self._literature_dto(item)
                data["project_count"] = project_counts.get(item.id, 0)
                data["record_count"] = record_counts.get(item.id, 0)
                result.append(data)
            if not payload.get("pagination"): return result
            return {"items": result, "pagination": {"page": page, "page_size": page_size, "pages": pages, "total": total}}
        return self._run(query)

    def get_literature(self, payload):
        item_id = _positive_id(payload.get("id"))

        def query(workspace):
            item = LiteratureItem.query.filter_by(id=item_id, workspace_id=workspace.id, is_deleted=False).first()
            if not item:
                raise NotFoundError("文献不存在或已移入回收站。")
            return self._literature_dto(item, detail=True)
        return self._run(query)

    def literature_facets(self, payload=None):
        def query(workspace):
            rows = db.session.execute(
                db.select(LiteratureItem.source, func.count()).where(
                    LiteratureItem.workspace_id == workspace.id, LiteratureItem.is_deleted.is_(False),
                ).group_by(LiteratureItem.source)
            ).all()
            sources = {key: count for key, count in rows}
            sources["all"] = sum(sources.values())
            return {"sources": sources}
        return self._run(query)

    def save_literature(self, payload, expected_row_version=None):
        item_id = _positive_id(payload.get("id"), optional=True)
        title = _string(payload.get("title"), 500, "title", required=True)
        source = _string(payload.get("source") or "manual", 20, "source")
        if source not in LITERATURE_SOURCES:
            raise ValidationError("文献来源无效。", field_errors={"source": "不支持的来源"})
        read_status = _string(payload.get("read_status") or "unread", 20, "read_status")
        if read_status not in READ_STATUSES:
            raise ValidationError("阅读状态无效。", field_errors={"read_status": "不支持的状态"})
        authors = payload.get("authors") or []
        if isinstance(authors, str):
            authors = [part.strip() for part in re.split(r"[;,，；]", authors) if part.strip()]
        if not isinstance(authors, list):
            raise ValidationError("作者格式无效。", field_errors={"authors": "应为作者列表"})
        keywords = payload.get("keywords") or []
        if isinstance(keywords, str):
            keywords = [part.strip() for part in re.split(r"[;,，；]", keywords) if part.strip()]
        year_value = payload.get("year")
        year = None
        if year_value not in (None, ""):
            try:
                year = int(year_value)
            except (TypeError, ValueError) as exc:
                raise ValidationError("年份无效。", field_errors={"year": "必须为年份"}) from exc
            if year < 1000 or year > date.today().year + 2:
                raise ValidationError("年份无效。", field_errors={"year": "年份超出合理范围"})

        def save(workspace):
            if item_id:
                item = LiteratureItem.query.filter_by(id=item_id, workspace_id=workspace.id, is_deleted=False).first()
                if not item:
                    raise NotFoundError("文献不存在或已移入回收站。")
                expected = expected_row_version if expected_row_version is not None else payload.get("row_version")
                try:
                    expected = int(expected)
                except (TypeError, ValueError) as exc:
                    raise ValidationError("缺少有效的文献版本。") from exc
                if item.row_version != expected:
                    raise ConflictError("文献已被更新，请重新加载后合并修改。")
                item.row_version += 1
            else:
                item = LiteratureItem(
                    workspace_id=workspace.id, source=source,
                    source_key=_string(payload.get("source_key"), 255, "source_key") or f"manual:{uuid.uuid4()}",
                )
                db.session.add(item)
            item.title = title
            item.source = source
            item.authors_json = _json([_string(value, 200, "authors") for value in authors[:100]])
            item.year = year
            item.journal = _string(payload.get("journal"), 255, "journal")
            item.volume = _string(payload.get("volume"), 40, "volume")
            item.issue = _string(payload.get("issue"), 40, "issue")
            item.pages = _string(payload.get("pages"), 80, "pages")
            item.doi = _string(payload.get("doi"), 255, "doi")
            item.doi_normalized = _normalize_doi(item.doi)
            item.url = _string(payload.get("url"), 1000, "url")
            item.abstract = _string(payload.get("abstract"), 100000, "abstract")
            item.keywords_json = _json([_string(value, 100, "keywords") for value in keywords[:100]])
            item.read_status = read_status
            item.reading_notes = _string(payload.get("reading_notes"), 100000, "reading_notes")
            db.session.flush()
            self._index_entity(workspace, "literature", item)
            self._activity(workspace, "literature_saved", f"保存文献：{item.title}", "literature_item", item.id)
            return self._literature_dto(item, detail=True)
        return self._run(save)

    def link_literature(self, payload):
        literature_id = _positive_id(payload.get("literature_id"), "literature_id")
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        record_id = _positive_id(payload.get("record_id"), "record_id", optional=True)
        if not project_id and not record_id:
            raise ValidationError("请选择项目或实验记录。")

        def link(workspace):
            item = LiteratureItem.query.filter_by(id=literature_id, workspace_id=workspace.id, is_deleted=False).first()
            if not item:
                raise NotFoundError("文献不存在。")
            if project_id:
                self._ensure_project(workspace.id, project_id)
                db.session.execute(project_literature.insert().prefix_with("OR IGNORE").values(
                    project_id=project_id, literature_id=item.id,
                    purpose=_string(payload.get("purpose") or "reference", 120, "purpose"), created_at=utcnow(),
                ))
                db.session.execute(
                    project_literature.update().where(
                        project_literature.c.project_id == project_id,
                        project_literature.c.literature_id == item.id,
                    ).values(purpose=_string(payload.get("purpose") or "reference", 120, "purpose"))
                )
            if record_id:
                record = LabRecord.query.filter_by(id=record_id, workspace_id=workspace.id, is_deleted=False).first()
                if not record:
                    raise NotFoundError("实验记录不存在。")
                db.session.execute(lab_record_literature.insert().prefix_with("OR IGNORE").values(
                    record_id=record_id, literature_id=item.id,
                    citation_role=_string(payload.get("citation_role") or "reference", 80, "citation_role"),
                    locator=_string(payload.get("locator"), 255, "locator"),
                    notes=_string(payload.get("notes"), 5000, "notes"), created_at=utcnow(),
                ))
                db.session.execute(
                    lab_record_literature.update().where(
                        lab_record_literature.c.record_id == record_id,
                        lab_record_literature.c.literature_id == item.id,
                    ).values(
                        citation_role=_string(payload.get("citation_role") or "reference", 80, "citation_role"),
                        locator=_string(payload.get("locator"), 255, "locator"),
                        notes=_string(payload.get("notes"), 5000, "notes"),
                    )
                )
            item.row_version += 1
            return self._literature_dto(item, detail=True)
        return self._run(link)

    def unlink_literature(self, payload):
        literature_id = _positive_id(payload.get("literature_id"), "literature_id")
        relation_type = _string(payload.get("relation_type"), 20, "relation_type", required=True)
        target_id = _positive_id(payload.get("target_id"), "target_id")
        if relation_type not in {"project", "record"}:
            raise ValidationError("仅支持解除项目或实验记录关系。")

        def unlink(workspace):
            item = LiteratureItem.query.filter_by(
                id=literature_id, workspace_id=workspace.id, is_deleted=False,
            ).first()
            if not item:
                raise NotFoundError("文献不存在。")
            if relation_type == "project":
                self._ensure_project(workspace.id, target_id)
                statement = project_literature.delete().where(
                    project_literature.c.project_id == target_id,
                    project_literature.c.literature_id == item.id,
                )
            else:
                record = LabRecord.query.filter_by(
                    id=target_id, workspace_id=workspace.id, is_deleted=False,
                ).first()
                if not record:
                    raise NotFoundError("实验记录不存在。")
                statement = lab_record_literature.delete().where(
                    lab_record_literature.c.record_id == target_id,
                    lab_record_literature.c.literature_id == item.id,
                )
            removed = db.session.execute(statement).rowcount
            if removed:
                item.row_version += 1
            return {"removed": bool(removed), "literature": self._literature_dto(item, detail=True)}
        return self._run(unlink)

    def literature_bulk(self, payload):
        raw_ids = payload.get("ids")
        if not isinstance(raw_ids, list) or not raw_ids:
            raise ValidationError("请选择至少一篇文献。")
        ids = list(dict.fromkeys(_positive_id(value, "ids") for value in raw_ids))[:500]
        action = _string(payload.get("action"), 30, "action", required=True)
        if action not in {"status", "project", "trash"}:
            raise ValidationError("不支持的文献批量操作。")

        def bulk(workspace):
            items = LiteratureItem.query.filter(
                LiteratureItem.workspace_id == workspace.id,
                LiteratureItem.id.in_(ids),
                LiteratureItem.is_deleted.is_(False),
            ).all()
            if len(items) != len(ids):
                raise NotFoundError("部分文献不存在或不属于当前工作区。")
            if action == "status":
                status = _string(payload.get("read_status") or payload.get("value"), 20, "read_status", required=True)
                if status not in READ_STATUSES:
                    raise ValidationError("阅读状态无效。")
                for item in items:
                    item.read_status = status
                    item.row_version += 1
                    self._index_entity(workspace, "literature", item)
            elif action == "project":
                project_id = _positive_id(payload.get("project_id") or payload.get("value"), "project_id")
                self._ensure_project(workspace.id, project_id)
                for item in items:
                    db.session.execute(project_literature.insert().prefix_with("OR IGNORE").values(
                        project_id=project_id, literature_id=item.id,
                        purpose=_string(payload.get("purpose") or "reference", 120, "purpose"),
                        created_at=utcnow(),
                    ))
                    item.row_version += 1
            else:
                for item in items:
                    item.is_deleted = True
                    item.deleted_at = utcnow()
                    item.row_version += 1
                    self._index_entity(workspace, "literature", item)
            return {"updated": len(items), "action": action, "results": [
                {"id": item.id, "status": "updated", "row_version": item.row_version} for item in items
            ]}
        return self._run(bulk)

    def zotero_status(self, payload=None):
        def query(workspace):
            connection = ZoteroConnection.query.filter_by(workspace_id=workspace.id).first()
            if not connection:
                connection = ZoteroConnection(workspace_id=workspace.id)
                db.session.add(connection)
                db.session.flush()
            return {
                "state": connection.connection_state, "base_url": connection.base_url,
                "last_sync_at": _iso(connection.last_sync_at), "last_success_at": _iso(connection.last_success_at),
                "last_error_code": connection.last_error_code, "last_error_message": connection.last_error_message,
                "server_id": connection.server_id, "library_version": connection.library_version,
                "library_key": connection.library_key or "personal",
                "sync_state": connection.sync_state, "sync_progress": connection.sync_progress,
                "sync_stage": connection.sync_stage,
            }
        return self._run(query)

    def zotero_sync(self, payload=None):
        payload = payload or {}

        def sync(workspace):
            connection = ZoteroConnection.query.filter_by(workspace_id=workspace.id).first()
            if not connection:
                connection = ZoteroConnection(workspace_id=workspace.id)
                db.session.add(connection)
                db.session.flush()
            base_url = _string(payload.get("base_url") or connection.base_url, 500, "base_url")
            if base_url.rstrip("/") != "http://127.0.0.1:23119":
                raise ValidationError("Zotero Local API 地址必须为 http://127.0.0.1:23119。")
            connection.base_url = base_url.rstrip("/")
            previous_library = connection.library_key or "personal"
            requested_library = _string(payload.get("library_key") or previous_library, 120, "library_key")
            library_changed = requested_library != previous_library
            if requested_library == "personal":
                library_path, source_prefix = "users/0", "zotero:0:"
            elif requested_library.startswith("group:") and requested_library[6:].isdigit():
                group_id = requested_library[6:]
                library_path, source_prefix = f"groups/{group_id}", f"zotero:groups:{group_id}:"
            else:
                raise ValidationError("Zotero 资料库必须为 personal 或 group:数字ID。")
            connection.library_key = requested_library
            connection.last_sync_at = utcnow()
            connection.sync_state = "running"
            connection.sync_progress = 5
            progress_callback = payload.get("_progress_callback")
            if callable(progress_callback):
                progress_callback(5, "读取 Zotero 变更")
            connection.sync_stage = "读取 Zotero 变更"
            try:
                force_full = bool(payload.get("force_full")) or library_changed or not connection.library_version
                since = None if force_full else connection.library_version
                cancel_event = payload.get("_cancel_event")
                items, metadata = _fetch_zotero_items(
                    connection.base_url, library_path, since=since, cancel_event=cancel_event,
                )
                server_changed = bool(
                    connection.server_id
                    and metadata["server_id"]
                    and connection.server_id != metadata["server_id"]
                )
                if server_changed and since is not None:
                    items, metadata = _fetch_zotero_items(
                        connection.base_url, library_path, since=None, cancel_event=cancel_event,
                    )
                    force_full = True
                complete_library = force_full and metadata["complete_library"]
                if complete_library and any(
                    not isinstance(raw, dict) or not isinstance(raw.get("data"), dict)
                    or not str(raw.get("key") or raw.get("data", {}).get("key") or "").strip()
                    for raw in items
                ):
                    raise ValueError("Zotero full response contains a malformed item")
                added = updated = restored = attachments = 0
                seen = set()
                missing = 0
                attachment_rows = []
                connection.sync_progress = 35
                if callable(progress_callback):
                    progress_callback(35, "更新文献与附件")
                connection.sync_stage = "更新文献与附件"
                with db.session.begin_nested():
                    staged_missing = []
                    for raw in items:
                        if cancel_event is not None and cancel_event.is_set():
                            raise ZoteroSyncCancelled()
                        if not isinstance(raw, dict):
                            continue
                        data = raw.get("data")
                        key = str(
                            raw.get("key")
                            or (data.get("key") if isinstance(data, dict) else "")
                            or ""
                        ).strip()
                        source_key = f"{source_prefix}{key}" if key else ""
                        if key:
                            source_key = f"{source_prefix}{key}"
                        if not isinstance(data, dict):
                            continue
                        item_type = str(data.get("itemType") or "")
                        if item_type == "attachment":
                            attachment_rows.append((raw, data, key, source_key))
                            continue
                        if item_type in {"note", "annotation"} or data.get("parentItem"):
                            continue
                        if key and complete_library:
                            seen.add(source_key)
                        if data.get("deleted"):
                            cached = LiteratureItem.query.filter_by(
                                workspace_id=workspace.id, source_key=source_key,
                            ).first()
                            if cached:
                                staged_missing.append(cached)
                            continue
                        title = str(data.get("title") or "").strip()
                        if not key or not title:
                            continue
                        item = LiteratureItem.query.filter_by(
                            workspace_id=workspace.id,
                            source_key=source_key,
                        ).first()
                        is_new = item is None
                        if is_new:
                            item = LiteratureItem(
                                workspace_id=workspace.id,
                                source="zotero",
                                source_key=source_key,
                                title=title,
                            )
                            db.session.add(item)
                            added += 1
                        else:
                            updated += 1
                            item.row_version += 1
                            if item.is_deleted:
                                item.is_deleted, item.deleted_at = False, None
                                restored += 1
                        creators = data.get("creators") if isinstance(data.get("creators"), list) else []
                        item.title = title
                        item.authors_json = _json([
                            " ".join(filter(None, [creator.get("firstName"), creator.get("lastName")])).strip()
                            or str(creator.get("name") or "")
                            for creator in creators
                            if isinstance(creator, dict)
                        ])
                        raw_date = str(data.get("date") or "")
                        match = re.search(r"(?:19|20)\d{2}", raw_date)
                        item.year = int(match.group()) if match else None
                        item.journal = str(
                            data.get("publicationTitle") or data.get("proceedingsTitle") or ""
                        )[:255]
                        item.volume = str(data.get("volume") or "")[:40]
                        item.issue = str(data.get("issue") or "")[:40]
                        item.pages = str(data.get("pages") or "")[:80]
                        item.doi = str(data.get("DOI") or "")[:255]
                        item.doi_normalized = _normalize_doi(item.doi)
                        item.url = str(data.get("url") or "")[:1000]
                        item.abstract = str(data.get("abstractNote") or "")
                        item.keywords_json = _json([
                            str(tag.get("tag"))
                            for tag in data.get("tags", [])
                            if isinstance(tag, dict) and tag.get("tag")
                        ])
                        try:
                            item.zotero_version = int(raw.get("version") or 0)
                        except (TypeError, ValueError):
                            item.zotero_version = 0
                        item.synced_at = utcnow()
                        item.source_missing = False
                        item.source_snapshot_json = _json(raw)
                    db.session.flush()
                    for raw, data, key, attachment_source_key in attachment_rows:
                        if not key:
                            continue
                        parent_key = str(data.get("parentItem") or "").strip()
                        parent = LiteratureItem.query.filter_by(
                            workspace_id=workspace.id,
                            source_key=f"{source_prefix}{parent_key}",
                        ).first() if parent_key else None
                        file_item = LibraryItem.query.filter_by(
                            workspace_id=workspace.id,
                            source_key=attachment_source_key,
                        ).first()
                        if not file_item:
                            file_item = LibraryItem(
                                workspace_id=workspace.id,
                                source_key=attachment_source_key,
                                storage_mode="zotero",
                                display_name=str(data.get("filename") or data.get("title") or key)[:255],
                            )
                            db.session.add(file_item)
                            attachments += 1
                        file_item.original_name = str(data.get("filename") or "")[:255]
                        file_item.description = str(data.get("title") or "")[:10000]
                        file_item.kind = "literature"
                        file_item.mime_type = str(data.get("contentType") or "application/octet-stream")[:160]
                        file_item.link_status = "missing" if data.get("deleted") else "unchecked"
                        file_item.ai_readability = "metadata_only"
                        file_item.external_path = None
                        file_item.managed_relative_path = None
                        file_item.path_normalized = ""
                        db.session.flush()
                        if parent:
                            db.session.execute(
                                literature_library_item.insert().prefix_with("OR IGNORE").values(
                                    literature_id=parent.id,
                                    library_item_id=file_item.id,
                                    attachment_role="full_text",
                                    created_at=utcnow(),
                                )
                            )
                    if complete_library:
                        for cached in LiteratureItem.query.filter_by(
                            workspace_id=workspace.id,
                            source="zotero",
                        ).filter(LiteratureItem.source_key.like(f"{source_prefix}%")).all():
                            if cached.source_key not in seen:
                                staged_missing.append(cached)
                    if cancel_event is not None and cancel_event.is_set():
                        raise ZoteroSyncCancelled()
                    for cached in {item.id: item for item in staged_missing}.values():
                        if not cached.source_missing:
                            missing += 1
                        cached.source_missing = True
                connection.connection_state = "connected"
                connection.last_success_at = utcnow()
                connection.last_error_code = ""
                connection.last_error_message = ""
                connection.server_id = metadata["server_id"] or connection.server_id
                connection.library_version = max(
                    connection.library_version if not (server_changed or library_changed) else 0,
                    metadata["library_version"],
                )
                if complete_library:
                    connection.last_full_sync_at = utcnow()
                else:
                    connection.last_incremental_sync_at = utcnow()
                connection.sync_state = "completed"
                connection.sync_progress = 100
                if callable(progress_callback):
                    progress_callback(100, "同步完成")
                connection.sync_stage = "同步完成"
                return {
                    "state": "connected",
                    "added": added,
                    "updated": updated,
                    "restored": restored,
                    "missing": missing,
                    "attachments": attachments,
                    "incremental": not complete_library,
                    "server_changed": server_changed,
                    "library_version": connection.library_version,
                }
            except ZoteroSyncCancelled:
                connection.sync_state = "cancelled"
                connection.sync_progress = 0
                connection.sync_stage = "同步已取消"
                connection.last_error_code = ""
                connection.last_error_message = ""
                return {"state": "cancelled", "added": 0, "updated": 0, "cancelled": True}
            except urllib.error.HTTPError as exc:
                connection.connection_state = "unavailable" if exc.code == 403 else "error"
                connection.sync_state = "failed"
                connection.sync_stage = "同步失败"
                connection.last_error_code = (
                    "zotero_api_disabled" if exc.code == 403 else "zotero_http_error"
                )
                connection.last_error_message = f"HTTP {exc.code}"[:500]
                message = (
                    "Zotero Local API 未启用或拒绝访问。"
                    if exc.code == 403
                    else f"Zotero Local API 返回 HTTP {exc.code}。"
                )
                return {
                    "state": connection.connection_state,
                    "added": 0,
                    "updated": 0,
                    "error": message,
                }
            except urllib.error.URLError as exc:
                connection.connection_state = "unavailable"
                connection.sync_state = "failed"
                connection.sync_stage = "同步失败"
                is_timeout = isinstance(exc.reason, TimeoutError)
                connection.last_error_code = "zotero_timeout" if is_timeout else "zotero_unavailable"
                connection.last_error_message = str(exc.reason)[:500]
                message = "连接 Zotero 超时。" if is_timeout else "Zotero 未运行或 Local API 不可用。"
                return {"state": "unavailable", "added": 0, "updated": 0, "error": message}
            except TimeoutError as exc:
                connection.connection_state = "unavailable"
                connection.sync_state = "failed"
                connection.sync_stage = "同步失败"
                connection.last_error_code = "zotero_timeout"
                connection.last_error_message = str(exc)[:500]
                return {"state": "unavailable", "added": 0, "updated": 0, "error": "连接 Zotero 超时。"}
            except (UnicodeDecodeError, json.JSONDecodeError, ValueError) as exc:
                connection.connection_state = "error"
                connection.sync_state = "failed"
                connection.sync_stage = "同步失败"
                connection.last_error_code = "invalid_response"
                connection.last_error_message = str(exc)[:500]
                return {"state": "error", "added": 0, "updated": 0, "error": "Zotero 返回了无法识别的数据。"}
        return self._run(sync)

    def zotero_collections_sync(self, payload=None):
        payload = payload or {}

        def sync(workspace):
            connection = ZoteroConnection.query.filter_by(workspace_id=workspace.id).first()
            if not connection:
                raise ValidationError("请先连接并同步 Zotero。")
            library_key = connection.library_key or "personal"
            if library_key == "personal":
                library_path, source_prefix = "users/0", "zotero:0:"
            elif library_key.startswith("group:") and library_key[6:].isdigit():
                group_id = library_key[6:]
                library_path, source_prefix = f"groups/{group_id}", f"zotero:groups:{group_id}:"
            else:
                raise ValidationError("Zotero 资料库标识无效。")
            rows, _ = _fetch_zotero_items(
                connection.base_url.rstrip("/"), library_path, resource="collections",
                cancel_event=payload.get("_cancel_event"),
            )
            seen = set()
            for raw in rows:
                data = raw.get("data", {}) if isinstance(raw, dict) else {}
                key = str(raw.get("key") or data.get("key") or "").strip() if isinstance(raw, dict) else ""
                if not key: continue
                seen.add(key)
                item = ZoteroCollection.query.filter_by(
                    workspace_id=workspace.id, library_key=library_key, collection_key=key,
                ).first()
                if not item:
                    item = ZoteroCollection(workspace_id=workspace.id, library_key=library_key, collection_key=key)
                    db.session.add(item)
                item.name = str(data.get("name") or key)[:500]
                item.parent_key = str(data.get("parentCollection") or "")[:120]
                try: item.version = int(raw.get("version") or data.get("version") or 0)
                except (TypeError, ValueError): item.version = 0
                item.source_missing = bool(data.get("deleted"))
            for cached in ZoteroCollection.query.filter_by(workspace_id=workspace.id, library_key=library_key).all():
                if cached.collection_key not in seen: cached.source_missing = True
            db.session.flush()
            collections = {item.collection_key: item for item in ZoteroCollection.query.filter_by(
                workspace_id=workspace.id, library_key=library_key, source_missing=False,
            ).all()}
            literature = LiteratureItem.query.filter_by(workspace_id=workspace.id, source="zotero").filter(
                LiteratureItem.source_key.like(f"{source_prefix}%")
            ).all()
            literature_ids = [item.id for item in literature]
            if literature_ids:
                db.session.execute(zotero_collection_literature.delete().where(
                    zotero_collection_literature.c.literature_id.in_(literature_ids)
                ))
            memberships = 0
            for item in literature:
                snapshot = _load_json(item.source_snapshot_json, {})
                data = snapshot.get("data", {}) if isinstance(snapshot, dict) else {}
                for key in data.get("collections", []) if isinstance(data.get("collections"), list) else []:
                    collection = collections.get(str(key))
                    if collection:
                        db.session.execute(zotero_collection_literature.insert().values(
                            collection_id=collection.id, literature_id=item.id,
                        ))
                        memberships += 1
            return {"collections": len(collections), "memberships": memberships}
        return self._run(sync)

    def zotero_collections(self, payload=None):
        payload = payload or {}

        def query(workspace):
            rows = ZoteroCollection.query.filter_by(workspace_id=workspace.id, source_missing=False).order_by(
                ZoteroCollection.parent_key, ZoteroCollection.name
            ).all()
            return [{
                "id": item.id, "key": item.collection_key, "name": item.name,
                "parent_key": item.parent_key, "library_key": item.library_key,
                "project_ids": [row[0] for row in db.session.execute(
                    db.select(project_zotero_collection.c.project_id).where(
                        project_zotero_collection.c.collection_id == item.id
                    )
                )],
            } for item in rows]
        return self._run(query)

    def zotero_collection_map(self, payload):
        collection_id = _positive_id(payload.get("collection_id"), "collection_id")
        project_id = _positive_id(payload.get("project_id"), "project_id")
        enabled = bool(payload.get("enabled", True))

        def update_mapping(workspace):
            collection = ZoteroCollection.query.filter_by(id=collection_id, workspace_id=workspace.id).first()
            if not collection: raise NotFoundError("Zotero Collection 不存在。")
            self._ensure_project(workspace.id, project_id)
            condition = (project_zotero_collection.c.project_id == project_id) & (project_zotero_collection.c.collection_id == collection_id)
            if enabled:
                db.session.execute(project_zotero_collection.insert().prefix_with("OR IGNORE").values(
                    project_id=project_id, collection_id=collection_id,
                ))
            else:
                db.session.execute(project_zotero_collection.delete().where(condition))
            return {"collection_id": collection_id, "project_id": project_id, "enabled": enabled}
        return self._run(update_mapping)

    @staticmethod
    def _library_dto(item):
        return {
            "id": item.id, "display_name": item.display_name, "original_name": item.original_name,
            "description": item.description, "kind": item.kind, "storage_mode": item.storage_mode,
            "path": item.managed_relative_path or item.external_path or "", "mime_type": item.mime_type,
            "size_bytes": item.size_bytes, "sha256": item.sha256, "link_status": item.link_status,
            "ai_readability": item.ai_readability, "updated_at": _iso(item.updated_at),
        }

    def list_library_items(self, payload=None):
        payload = payload or {}
        search = _string(payload.get("search"), 200, "search")
        kind = _string(payload.get("kind"), 30, "kind")
        mode = _string(payload.get("storage_mode"), 20, "storage_mode")
        record_id = _positive_id(payload.get("record_id"), "record_id", optional=True)
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)

        def query(workspace):
            items = LibraryItem.query.filter_by(workspace_id=workspace.id, is_deleted=False)
            if record_id:
                items = items.join(
                    lab_record_library_item,
                    lab_record_library_item.c.library_item_id == LibraryItem.id,
                ).filter(lab_record_library_item.c.record_id == record_id)
            if project_id:
                self._ensure_project(workspace.id, project_id)
                items = items.join(
                    project_library_item,
                    project_library_item.c.library_item_id == LibraryItem.id,
                ).filter(project_library_item.c.project_id == project_id)
            if kind:
                items = items.filter_by(kind=kind)
            if mode:
                items = items.filter_by(storage_mode=mode)
            if search:
                pattern = f"%{search}%"
                items = items.filter(or_(LibraryItem.display_name.ilike(pattern), LibraryItem.external_path.ilike(pattern)))
            items = _sorted(items, payload, {
                "updated_desc": (LibraryItem.updated_at.desc(),),
                "name_asc": (LibraryItem.display_name.asc(),),
                "size_desc": (LibraryItem.size_bytes.desc(), LibraryItem.display_name.asc()),
            }, "updated_desc")
            return _page_result(items, payload, self._library_dto)
        return self._run(query)

    def library_item_path(self, payload):
        item_id = _positive_id(payload.get("id"))

        def resolve(workspace):
            item = LibraryItem.query.filter_by(id=item_id, workspace_id=workspace.id, is_deleted=False).first()
            if not item:
                raise NotFoundError("文件不存在或已移入回收站。")
            if item.storage_mode != "zotero":
                raw = item.managed_relative_path or item.external_path
                path = ((Path(self.flask_app.instance_path) / "library" / raw).resolve()
                        if item.managed_relative_path else Path(raw or "").expanduser().resolve())
                if not path.is_file():
                    raise NotFoundError("文件已失联，请重新定位或校验文件。")
                return {"path": str(path), "storage_mode": item.storage_mode}
            source_key = item.source_key or ""
            personal = re.fullmatch(r"zotero:0:([A-Za-z0-9]+)", source_key)
            group = re.fullmatch(r"zotero:groups:(\d+):([A-Za-z0-9]+)", source_key)
            if personal:
                library_path, key = "users/0", personal.group(1)
            elif group:
                library_path, key = f"groups/{group.group(1)}", group.group(2)
            else:
                raise ValidationError("Zotero 附件标识无效。")
            connection = ZoteroConnection.query.filter_by(workspace_id=workspace.id).first()
            base_url = (connection.base_url if connection else "http://127.0.0.1:23119").rstrip("/")
            request = urllib.request.Request(
                f"{base_url}/api/{library_path}/items/{key}/file/view/url",
                headers={"Accept": "text/plain", "Zotero-API-Version": "3"}, method="GET",
            )
            try:
                with urllib.request.urlopen(request, timeout=8) as response:
                    raw_url = response.read().decode("utf-8").strip().strip('"')
            except (urllib.error.URLError, urllib.error.HTTPError, TimeoutError) as exc:
                raise NotFoundError("无法从 Zotero 获取附件，请启动 Zotero 或修复附件后重试。") from exc
            parsed = urllib.parse.urlparse(raw_url)
            if parsed.scheme != "file":
                raise ValidationError("Zotero 返回了不安全的附件位置。")
            file_path = urllib.request.url2pathname(urllib.parse.unquote(parsed.path))
            if parsed.netloc:
                file_path = f"//{parsed.netloc}{file_path}"
            path = Path(file_path).resolve()
            if not path.is_file():
                raise NotFoundError("Zotero 附件不存在，请在 Zotero 中重新定位文件。")
            return {"path": str(path), "storage_mode": "zotero"}
        return self._run(resolve)

    def import_library_item(self, payload):
        source_path = Path(_string(payload.get("path"), 2000, "path", required=True)).expanduser().resolve()
        if not source_path.is_file():
            raise ValidationError("所选文件不存在或不可读取。", field_errors={"path": "文件无效"})
        storage_mode = _string(payload.get("storage_mode") or "managed", 20, "storage_mode")
        if storage_mode not in STORAGE_MODES:
            raise ValidationError("存储方式无效。")
        kind = _string(payload.get("kind") or "other", 30, "kind")
        if kind not in FILE_KINDS:
            raise ValidationError("文件分类无效。")
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        record_id = _positive_id(payload.get("record_id"), "record_id", optional=True)

        with self.flask_app.app_context():
            operation = None
            target_path = None
            try:
                workspace = self._workspace()
                self._ensure_project(workspace.id, project_id)
                if record_id:
                    record = LabRecord.query.filter_by(
                        id=record_id, workspace_id=workspace.id, is_deleted=False
                    ).first()
                    if not record:
                        raise NotFoundError("关联实验记录不存在。")
                    if project_id and record.project_id != project_id:
                        raise ValidationError("文件的项目与实验记录不一致。")
                mime_type = mimetypes.guess_type(source_path.name)[0] or "application/octet-stream"
                file_hash = _sha256(source_path)
                size = source_path.stat().st_size
                if storage_mode == "managed":
                    root = (Path(self.flask_app.instance_path) / "library").resolve()
                    root.mkdir(parents=True, exist_ok=True)
                    safe_name = re.sub(r"[^\w.()\- \u4e00-\u9fff]", "_", source_path.name)[:180] or "file"
                    relative = Path(date.today().strftime("%Y/%m")) / f"{uuid.uuid4().hex[:12]}-{safe_name}"
                    target_path = (root / relative).resolve()
                    if root != target_path and root not in target_path.parents:
                        raise ValidationError("托管文件路径越界。")
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                    operation = FileOperation(
                        operation_type="copy", source_path=str(source_path), target_path=str(target_path), status="prepared",
                    )
                    db.session.add(operation)
                    db.session.commit()
                    shutil.copy2(source_path, target_path)
                    operation.status = "file_done"
                    db.session.commit()
                    managed_relative_path = relative.as_posix()
                    external_path = None
                    path_normalized = managed_relative_path.casefold()
                else:
                    managed_relative_path = None
                    external_path = str(source_path)
                    path_normalized = str(source_path).replace("/", "\\").casefold()
                item = LibraryItem(
                    workspace_id=workspace.id,
                    display_name=_string(payload.get("display_name") or source_path.name, 255, "display_name", required=True),
                    original_name=source_path.name, description=_string(payload.get("description"), 20000, "description"),
                    kind=kind, storage_mode=storage_mode, managed_relative_path=managed_relative_path,
                    external_path=external_path, path_normalized=path_normalized, mime_type=mime_type,
                    size_bytes=size, sha256=file_hash, link_status="available", last_verified_at=utcnow(),
                    ai_readability="image" if mime_type.startswith("image/") else "text" if mime_type.startswith("text/") else "metadata_only",
                )
                db.session.add(item)
                db.session.flush()
                self._index_entity(workspace, "file", item)
                if project_id:
                    db.session.execute(project_library_item.insert().values(
                        project_id=project_id, library_item_id=item.id, relation_role="evidence", created_at=utcnow(),
                    ))
                if record_id:
                    db.session.execute(lab_record_library_item.insert().values(
                        record_id=record_id, library_item_id=item.id, section_key="files",
                        relation_role="evidence", created_at=utcnow(),
                    ))
                if operation:
                    operation.library_item_id = item.id
                    operation.status = "db_done"
                self._activity(workspace, "file_imported", f"导入文件：{item.display_name}", "library_item", item.id, project_id=project_id)
                db.session.commit()
                return {**self._library_dto(item), "duplicate_count": LibraryItem.query.filter(
                    LibraryItem.workspace_id == workspace.id, LibraryItem.sha256 == file_hash,
                    LibraryItem.size_bytes == size, LibraryItem.id != item.id, LibraryItem.is_deleted.is_(False),
                ).count()}
            except Exception as exc:
                db.session.rollback()
                if operation:
                    operation.status = "failed"
                    operation.error_message = str(exc)[:1000]
                    db.session.add(operation)
                    db.session.commit()
                raise
            finally:
                db.session.remove()

    def verify_library_item(self, payload):
        item_id = _positive_id(payload.get("id"))

        def verify(workspace):
            item = LibraryItem.query.filter_by(id=item_id, workspace_id=workspace.id, is_deleted=False).first()
            if not item:
                raise NotFoundError("文件索引不存在。")
            if item.storage_mode == "zotero":
                item.link_status = "unchecked"
                item.last_verified_at = utcnow()
                return {**self._library_dto(item), "verification_message": "Zotero 附件将在打开时通过 Local API 校验。"}
            if item.storage_mode == "managed":
                path = (Path(self.flask_app.instance_path) / "library" / item.managed_relative_path).resolve()
            else:
                path = Path(item.external_path or "")
            if not path.is_file():
                item.link_status = "missing"
            else:
                current_hash = _sha256(path)
                item.link_status = "available" if current_hash == item.sha256 else "changed"
                item.size_bytes = path.stat().st_size
            item.last_verified_at = utcnow()
            return self._library_dto(item)
        return self._run(verify)

    @staticmethod
    def _note_dto(item, *, detail=False):
        data = {
            "id": item.id, "project_id": item.project_id, "title": item.title, "kind": item.kind,
            "row_version": item.row_version, "updated_at": _iso(item.updated_at),
            "excerpt": re.sub(r"\s+", " ", item.body)[:140],
        }
        if detail:
            data.update({
                "body": item.body,
                "record_ids": [row[0] for row in db.session.execute(
                    db.select(note_lab_record.c.record_id).where(note_lab_record.c.note_id == item.id)
                )],
                "literature_ids": [row[0] for row in db.session.execute(
                    db.select(note_literature.c.literature_id).where(note_literature.c.note_id == item.id)
                )],
            })
        return data

    def list_notes(self, payload=None):
        payload = payload or {}
        search = _string(payload.get("search"), 200, "search")
        kind = _string(payload.get("kind"), 30, "kind")
        project_filter = payload.get("project_id")

        def query(workspace):
            items = Note.query.filter_by(workspace_id=workspace.id, is_deleted=False)
            if kind:
                items = items.filter_by(kind=kind)
            if project_filter == "unclassified":
                items = items.filter(Note.project_id.is_(None))
            elif project_filter not in (None, ""):
                project_id = _positive_id(project_filter, "project_id")
                self._ensure_project(workspace.id, project_id)
                items = items.filter(Note.project_id == project_id)
            if search:
                pattern = f"%{search}%"
                items = items.filter(or_(Note.title.ilike(pattern), Note.body.ilike(pattern)))
            items = _sorted(items, payload, {
                "updated_desc": (Note.updated_at.desc(),), "title_asc": (Note.title.asc(),),
            }, "updated_desc")
            return _page_result(items, payload, self._note_dto)
        return self._run(query)

    def note_bulk(self, payload):
        ids = sorted({_positive_id(value) for value in (payload.get("ids") or [])})
        if not ids or len(ids) > 500:
            raise ValidationError("请选择 1 至 500 条笔记。")
        action = _string(payload.get("action"), 20, "action", required=True)
        if action not in {"project", "kind", "trash"}:
            raise ValidationError("不支持的笔记批量操作。")

        def bulk(workspace):
            items = Note.query.filter(Note.workspace_id == workspace.id, Note.is_deleted.is_(False), Note.id.in_(ids)).all()
            if action == "project":
                value = _positive_id(payload.get("value"), "value", optional=True)
                self._ensure_project(workspace.id, value)
            elif action == "kind":
                value = _string(payload.get("value"), 30, "value", required=True)
                if value not in NOTE_KINDS:
                    raise ValidationError("笔记类型无效。")
            else:
                value = None
            results = []
            for item in items:
                if action == "trash": item.is_deleted, item.deleted_at = True, utcnow()
                elif action == "project": item.project_id = value
                else: item.kind = value
                item.row_version += 1
                self._index_entity(workspace, "note", item)
                results.append({"id": item.id, "status": "updated"})
            found = {item.id for item in items}
            results.extend({"id": value, "status": "not_found"} for value in ids if value not in found)
            return {"updated": len(items), "skipped": len(ids) - len(items), "results": results}
        return self._run(bulk)

    def get_note(self, payload):
        note_id = _positive_id(payload.get("id"))

        def query(workspace):
            item = Note.query.filter_by(id=note_id, workspace_id=workspace.id, is_deleted=False).first()
            if not item:
                raise NotFoundError("笔记不存在或已移入回收站。")
            return self._note_dto(item, detail=True)
        return self._run(query)

    def save_note(self, payload, expected_row_version=None):
        note_id = _positive_id(payload.get("id"), optional=True)
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        title = _string(payload.get("title"), 240, "title", required=True)
        kind = _string(payload.get("kind") or "general", 30, "kind")
        if kind not in NOTE_KINDS:
            raise ValidationError("笔记类型无效。")
        body = _string(payload.get("body"), 200000, "body")

        def save(workspace):
            self._ensure_project(workspace.id, project_id)
            if note_id:
                item = Note.query.filter_by(id=note_id, workspace_id=workspace.id, is_deleted=False).first()
                if not item:
                    raise NotFoundError("笔记不存在。")
                try:
                    expected = int(expected_row_version)
                except (TypeError, ValueError) as exc:
                    raise ValidationError("缺少有效的笔记版本。") from exc
                if item.row_version != expected:
                    raise ConflictError("笔记已被更新，请重新加载后合并修改。")
                item.row_version += 1
            else:
                item = Note(workspace_id=workspace.id, row_version=1)
                db.session.add(item)
            item.title, item.kind, item.body, item.project_id = title, kind, body, project_id
            db.session.flush()
            self._index_entity(workspace, "note", item)
            self._activity(workspace, "note_saved", f"保存笔记：{item.title}", "note", item.id, project_id=project_id)
            return self._note_dto(item, detail=True)
        return self._run(save)

    @staticmethod
    def _task_dto(item):
        status = {"待办": "todo", "进行中": "doing", "已完成": "done"}.get(item.status, item.status)
        priority = {"低": "low", "中": "medium", "高": "high"}.get(item.priority, item.priority)
        return {
            "id": item.id, "project_id": item.project_id, "lab_record_id": item.lab_record_id,
            "title": item.title, "category": item.category, "priority": priority,
            "deadline": _iso(item.deadline), "status": status, "notes": item.notes or "",
            "completed_at": _iso(item.completed_at), "row_version": item.row_version,
            "updated_at": _iso(item.updated_at),
        }

    def list_tasks(self, payload=None):
        payload = payload or {}
        scope = _string(payload.get("scope") or "all", 20, "scope")
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        search = _string(payload.get("search"), 160, "search")
        status_filter = _string(payload.get("status"), 20, "status")
        date_start = _date(payload.get("date_start"), "date_start")
        date_end = _date(payload.get("date_end"), "date_end")
        if date_start and date_end and date_end < date_start:
            raise ValidationError("结束日期不能早于开始日期。")
        today = date.today()

        def query(workspace):
            items = Task.query.filter_by(workspace_id=workspace.id, is_deleted=False)
            if project_id:
                self._ensure_project(workspace.id, project_id)
                items = items.filter(Task.project_id == project_id)
            if search:
                escaped = search.replace("%", r"\%").replace("_", r"\_")
                items = items.filter(or_(Task.title.ilike(f"%{escaped}%", escape="\\"), Task.notes.ilike(f"%{escaped}%", escape="\\")))
            if status_filter:
                items = items.filter(Task.status == status_filter)
            if date_start:
                items = items.filter(Task.deadline >= date_start)
            if date_end:
                items = items.filter(Task.deadline <= date_end)
            if scope == "today":
                items = items.filter(Task.deadline <= today, Task.status.notin_(("done", "cancelled", "已完成")))
            elif scope == "week":
                items = items.filter(Task.deadline.between(today, today + timedelta(days=7)))
            elif scope == "future":
                items = items.filter(Task.deadline > today + timedelta(days=7))
            elif scope == "completed":
                items = items.filter(Task.status.in_(("done", "cancelled", "已完成")))
            else:
                items = items.filter(Task.status.notin_(("cancelled",)))
            items = _sorted(items, payload, {
                "deadline_asc": (Task.deadline.is_(None), Task.deadline, Task.position),
                "updated_desc": (Task.updated_at.desc(),), "title_asc": (Task.title.asc(),),
            }, "deadline_asc")
            return _page_result(items, payload, self._task_dto)
        return self._run(query)

    def task_bulk(self, payload):
        ids = sorted({_positive_id(value) for value in (payload.get("ids") or [])})
        if not ids or len(ids) > 500:
            raise ValidationError("请选择 1 至 500 个任务。")
        action = _string(payload.get("action"), 20, "action", required=True)
        if action not in {"status", "project", "deadline", "trash"}:
            raise ValidationError("不支持的任务批量操作。")

        def bulk(workspace):
            items = Task.query.filter(Task.workspace_id == workspace.id, Task.is_deleted.is_(False), Task.id.in_(ids)).all()
            if action == "status":
                value = _string(payload.get("value"), 20, "value", required=True)
                if value not in TASK_STATUSES: raise ValidationError("任务状态无效。")
            elif action == "project":
                value = _positive_id(payload.get("value"), "value", optional=True)
                self._ensure_project(workspace.id, value)
            elif action == "deadline":
                value = _date(payload.get("value"), "value")
            else: value = None
            results = []
            for item in items:
                if action == "trash": item.is_deleted, item.deleted_at = True, utcnow()
                elif action == "status": item.status = value
                elif action == "project": item.project_id = value
                else: item.deadline = value
                item.row_version += 1
                results.append({"id": item.id, "status": "updated"})
            found = {item.id for item in items}
            results.extend({"id": value, "status": "not_found"} for value in ids if value not in found)
            return {"updated": len(items), "skipped": len(ids) - len(items), "results": results}
        return self._run(bulk)

    def save_task(self, payload, expected_row_version=None):
        task_id = _positive_id(payload.get("id"), optional=True)
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        title = _string(payload.get("title"), 160, "title", required=True)
        status = _string(payload.get("status") or "todo", 20, "status")
        priority = _string(payload.get("priority") or "medium", 10, "priority")
        if status not in TASK_STATUSES or priority not in TASK_PRIORITIES:
            raise ValidationError("任务状态或优先级无效。")
        record_id = _positive_id(payload.get("lab_record_id"), "lab_record_id", optional=True)
        deadline = _date(payload.get("deadline"), "deadline")

        def save(workspace):
            self._ensure_project(workspace.id, project_id)
            if record_id:
                record = LabRecord.query.filter_by(id=record_id, workspace_id=workspace.id, is_deleted=False).first()
                if not record:
                    raise NotFoundError("关联实验记录不存在。")
                if project_id and record.project_id != project_id:
                    raise ValidationError("任务的项目与实验记录不一致。")
                project = record.project_id
            else:
                project = project_id
            if task_id:
                item = Task.query.filter_by(id=task_id, workspace_id=workspace.id, is_deleted=False).first()
                if not item:
                    raise NotFoundError("任务不存在。")
                if expected_row_version is None:
                    raise ValidationError("缺少任务版本。")
                if item.row_version != int(expected_row_version):
                    raise ConflictError("任务已被更新，请重新加载。")
                item.row_version += 1
            else:
                item = Task(user_id=workspace.legacy_user_id, workspace_id=workspace.id, row_version=1)
                db.session.add(item)
            item.title, item.status, item.priority = title, status, priority
            item.category = _string(payload.get("category") or "research", 40, "category")
            item.deadline, item.notes = deadline, _string(payload.get("notes"), 20000, "notes")
            item.project_id, item.lab_record_id = project, record_id
            item.completed_at = utcnow() if status in {"done", "已完成"} else None
            db.session.flush()
            self._index_entity(workspace, "task", item)
            self._activity(workspace, "task_saved", f"保存任务：{item.title}", "task", item.id, project_id=project, record_id=record_id)
            return self._task_dto(item)
        return self._run(save)

    def list_calendar(self, payload=None):
        payload = payload or {}
        start = _date(payload.get("start"), "start", default=date.today().replace(day=1))
        end = _date(payload.get("end"), "end", default=start + timedelta(days=42))
        if end < start or (end - start).days > 370:
            raise ValidationError("日历范围无效。")

        def query(workspace):
            events = []
            for record in LabRecord.query.filter(
                LabRecord.workspace_id == workspace.id, LabRecord.is_deleted.is_(False),
                LabRecord.experiment_date.between(start, end),
            ).all():
                events.append({"id": f"record:{record.id}", "source_type": "record", "source_id": record.id,
                               "title": record.title, "date": _iso(record.experiment_date), "event_type": "experiment",
                               "project_id": record.project_id, "lab_record_id": record.id, "movable": False})
            for task in Task.query.filter(
                Task.workspace_id == workspace.id, Task.is_deleted.is_(False), Task.deadline.between(start, end),
            ).all():
                events.append({"id": f"task:{task.id}", "source_type": "task", "source_id": task.id,
                               "title": task.title, "date": _iso(task.deadline), "event_type": "task",
                               "project_id": task.project_id, "lab_record_id": task.lab_record_id, "movable": True})
            for item in CalendarEvent.query.filter(
                CalendarEvent.workspace_id == workspace.id, CalendarEvent.is_deleted.is_(False),
                CalendarEvent.starts_at >= datetime.combine(start, time.min),
                CalendarEvent.starts_at <= datetime.combine(end, time.max),
            ).all():
                events.append({"id": f"event:{item.id}", "source_type": "event", "source_id": item.id,
                               "title": item.title, "date": _iso(item.starts_at), "event_type": item.event_type,
                               "starts_at": _iso(item.starts_at), "ends_at": _iso(item.ends_at),
                               "all_day": item.all_day, "notes": item.notes, "project_id": item.project_id,
                               "lab_record_id": item.lab_record_id, "row_version": item.row_version,
                               "movable": True})
            for report in WeeklyReport.query.filter(
                WeeklyReport.workspace_id == workspace.id, WeeklyReport.is_deleted.is_(False),
                WeeklyReport.report_date.between(start, end),
            ).all():
                events.append({"id": f"weekly:{report.id}", "source_type": "weekly", "source_id": report.id,
                               "title": report.title, "date": _iso(report.report_date), "event_type": "weekly",
                               "project_id": report.project_id, "movable": False})
            return sorted(events, key=lambda item: item["date"] or "")
        return self._run(query)

    def create_calendar_event(self, payload):
        title = _string(payload.get("title"), 240, "title", required=True)
        starts_at = _datetime(payload.get("starts_at"), "starts_at")
        ends_at = _datetime(payload.get("ends_at"), "ends_at") if payload.get("ends_at") else None
        if ends_at and ends_at < starts_at:
            raise ValidationError("结束时间不能早于开始时间。")

        def create(workspace):
            item = CalendarEvent(
                workspace_id=workspace.id, title=title,
                event_type=_string(payload.get("event_type") or "meeting", 30, "event_type"),
                starts_at=starts_at, ends_at=ends_at, all_day=bool(payload.get("all_day")),
                notes=_string(payload.get("notes"), 20000, "notes"),
                project_id=_positive_id(payload.get("project_id"), "project_id", optional=True),
                lab_record_id=_positive_id(payload.get("lab_record_id"), "lab_record_id", optional=True), row_version=1,
            )
            self._ensure_project(workspace.id, item.project_id)
            db.session.add(item)
            db.session.flush()
            return {"id": item.id, "title": item.title, "starts_at": _iso(item.starts_at),
                    "ends_at": _iso(item.ends_at), "event_type": item.event_type, "all_day": item.all_day,
                    "notes": item.notes, "project_id": item.project_id, "lab_record_id": item.lab_record_id,
                    "row_version": item.row_version}
        return self._run(create)

    @staticmethod
    def _week_bounds(value=None):
        target = _date(value, "week", default=date.today())
        start = target - timedelta(days=target.weekday())
        return start, start + timedelta(days=6)

    def weekly_current(self, payload=None):
        payload = payload or {}
        start, end = self._week_bounds(payload.get("week"))
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)

        def query(workspace):
            records = LabRecord.query.filter(
                LabRecord.workspace_id == workspace.id, LabRecord.is_deleted.is_(False),
                LabRecord.experiment_date.between(start, end),
            ).order_by(LabRecord.experiment_date).all()
            tasks = Task.query.filter(
                Task.workspace_id == workspace.id, Task.is_deleted.is_(False),
                Task.completed_at >= datetime.combine(start, time.min),
                Task.completed_at <= datetime.combine(end, time.max),
            ).order_by(Task.completed_at).all()
            literature = LiteratureItem.query.filter(
                LiteratureItem.workspace_id == workspace.id, LiteratureItem.is_deleted.is_(False),
                LiteratureItem.updated_at >= datetime.combine(start, time.min),
                LiteratureItem.updated_at <= datetime.combine(end, time.max),
            ).order_by(LiteratureItem.updated_at).all()
            notes = Note.query.filter(
                Note.workspace_id == workspace.id, Note.is_deleted.is_(False),
                Note.updated_at >= datetime.combine(start, time.min), Note.updated_at <= datetime.combine(end, time.max),
            ).order_by(Note.updated_at).all()
            report = WeeklyReport.query.filter_by(
                workspace_id=workspace.id, period_start=start, period_end=end, is_deleted=False,
            ).order_by(WeeklyReport.updated_at.desc()).first()
            entries = []
            for source_type, values, title_attr, excerpt_attr, date_attr in (
                ("record", records, "title", "conclusion", "experiment_date"),
                ("task", tasks, "title", "notes", "completed_at"),
                ("literature", literature, "title", "abstract", "updated_at"),
                ("note", notes, "title", "body", "updated_at"),
            ):
                for item in values:
                    source_date = getattr(item, date_attr)
                    if isinstance(source_date, datetime):
                        source_date = source_date.date()
                    entries.append({"source_type": source_type, "source_id": item.id,
                                    "source_title": getattr(item, title_attr),
                                    "source_excerpt": (getattr(item, excerpt_attr) or "")[:300],
                                    "source_date": _iso(source_date), "include_state": "included"})
            project_ids = {item.project_id for item in records if item.project_id}
            return {
                "id": report.id if report else None, "row_version": report.row_version if report else 0,
                "title": report.title if report else f"{start:%Y 年第 %W 周}研究周报",
                "period_start": _iso(start), "period_end": _iso(end),
                "body": report.body if report else "", "issues_and_feedback": report.issues_and_feedback if report else "",
                "next_week_plan": report.next_week_plan if report else "",
                "counts": {"records": len(records), "tasks": len(tasks), "literature": len(literature),
                           "notes": len(notes), "projects": len(project_ids)},
                "entries": entries,
            }
        return self._run(query)

    def save_weekly(self, payload, expected_row_version=None):
        report_id = _positive_id(payload.get("id"), optional=True)
        start, end = self._week_bounds(payload.get("period_start"))
        body = _string(payload.get("body"), 200000, "body")
        issues = _string(payload.get("issues_and_feedback"), 100000, "issues_and_feedback")
        next_plan = _string(payload.get("next_week_plan"), 100000, "next_week_plan")
        entries = payload.get("entries") or []
        if not isinstance(entries, list):
            raise ValidationError("周报来源格式无效。")

        def save(workspace):
            if report_id:
                report = WeeklyReport.query.filter_by(id=report_id, workspace_id=workspace.id, is_deleted=False).first()
                if not report:
                    raise NotFoundError("周报不存在。")
                if expected_row_version is None:
                    raise ValidationError("缺少周报版本。")
                if report.row_version != int(expected_row_version):
                    raise ConflictError("周报已被更新，请重新加载。")
                report.row_version += 1
                WeeklyReportEntry.query.filter_by(report_id=report.id).delete()
            else:
                report = WeeklyReport(
                    user_id=workspace.legacy_user_id, workspace_id=workspace.id,
                    title=_string(payload.get("title") or f"{start:%Y 年第 %W 周}研究周报", 180, "title"),
                    report_date=end, period_start=start, period_end=end, row_version=1,
                )
                db.session.add(report)
            report.body, report.issues_and_feedback, report.next_week_plan = body, issues, next_plan
            report.summary = body[:500]
            report.source_snapshot_json = _json({"period_start": _iso(start), "period_end": _iso(end), "entry_count": len(entries)})
            db.session.flush()
            for position, raw in enumerate(entries, start=1):
                if not isinstance(raw, dict) or raw.get("include_state") == "excluded":
                    continue
                db.session.add(WeeklyReportEntry(
                    report_id=report.id, position=position,
                    source_type=_string(raw.get("source_type"), 30, "source_type", required=True),
                    source_id=_positive_id(raw.get("source_id"), "source_id"),
                    source_title=_string(raw.get("source_title"), 500, "source_title", required=True),
                    source_excerpt=_string(raw.get("source_excerpt"), 5000, "source_excerpt"),
                    source_date=_date(raw.get("source_date"), "source_date"), include_state="included",
                ))
            db.session.flush()
            self._index_entity(workspace, "weekly", report)
            self._activity(workspace, "weekly_saved", f"保存周报：{report.title}", "weekly_report", report.id)
            return {"id": report.id, "row_version": report.row_version, "updated_at": _iso(report.updated_at)}
        return self._run(save)

    @staticmethod
    def _weekly_version_number(role):
        match = re.fullmatch(r"version:(\d+)", str(role or ""))
        return int(match.group(1)) if match else 0

    @classmethod
    def _weekly_file_dto(cls, item, role):
        return {
            "id": item.id,
            "version_number": cls._weekly_version_number(role),
            "display_name": item.display_name,
            "original_name": item.original_name,
            "storage_mode": item.storage_mode,
            "mime_type": item.mime_type,
            "size_bytes": item.size_bytes,
            "sha256": item.sha256,
            "link_status": item.link_status,
            "updated_at": _iso(item.updated_at),
        }

    @staticmethod
    def _weekly_file_rows(report_id):
        return db.session.query(
            LibraryItem, weekly_report_library_item.c.relation_role,
        ).join(
            weekly_report_library_item,
            weekly_report_library_item.c.library_item_id == LibraryItem.id,
        ).filter(
            weekly_report_library_item.c.report_id == report_id,
            LibraryItem.is_deleted.is_(False),
        ).all()

    def import_weekly_file(self, payload, expected_row_version=None):
        source_path = Path(_string(payload.get("path"), 2000, "path", required=True)).expanduser().resolve()
        if not source_path.is_file():
            raise ValidationError("所选周报文件不存在或不可读取。", field_errors={"path": "文件无效"})
        if source_path.suffix.lower() not in WEEKLY_FILE_EXTENSIONS:
            raise ValidationError("周报支持 Word、PDF、PPT、表格、文本和 ZIP 文件。")
        storage_mode = _string(payload.get("storage_mode") or "managed", 20, "storage_mode")
        if storage_mode not in {"managed", "external"}:
            raise ValidationError("周报文件仅支持托管副本或外部链接。")
        report_id = _positive_id(payload.get("report_id"), "report_id", optional=True)
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        report_date = _date(payload.get("report_date"), "report_date", default=date.today())
        period_start = _date(payload.get("period_start"), "period_start")
        period_end = _date(payload.get("period_end"), "period_end")
        if period_start and period_end and period_end < period_start:
            raise ValidationError("周报结束日期不能早于开始日期。")

        with self.flask_app.app_context():
            target_path = None
            temporary_path = None
            try:
                workspace = self._workspace()
                self._ensure_project(workspace.id, project_id)
                if report_id:
                    report = WeeklyReport.query.filter_by(
                        id=report_id, workspace_id=workspace.id, is_deleted=False,
                    ).first()
                    if not report:
                        raise NotFoundError("周报不存在或已移入回收站。")
                    if expected_row_version is None:
                        raise ValidationError("缺少周报版本。")
                    if report.row_version != int(expected_row_version):
                        raise ConflictError("周报已被更新，请重新加载。")
                    report.row_version += 1
                else:
                    report = WeeklyReport(
                        user_id=workspace.legacy_user_id,
                        workspace_id=workspace.id,
                        title=source_path.stem[:180],
                        report_date=report_date,
                        row_version=1,
                    )
                    db.session.add(report)
                    db.session.flush()

                existing_rows = self._weekly_file_rows(report.id)
                version_number = max(
                    (self._weekly_version_number(role) for _, role in existing_rows),
                    default=0,
                ) + 1
                safe_name = re.sub(r"[^\w.()\- \u4e00-\u9fff]", "_", source_path.name)[:180] or "weekly-report"
                mime_type = mimetypes.guess_type(source_path.name)[0] or "application/octet-stream"
                file_hash = _sha256(source_path)
                size_bytes = source_path.stat().st_size
                if size_bytes <= 0:
                    raise ValidationError("不能导入空的周报文件。")

                if storage_mode == "managed":
                    root = (Path(self.flask_app.instance_path) / "library").resolve()
                    relative = Path("weekly") / f"report-{report.id}" / f"v{version_number:04d}" / f"{uuid.uuid4().hex[:12]}-{safe_name}"
                    target_path = (root / relative).resolve()
                    if root != target_path and root not in target_path.parents:
                        raise ValidationError("托管周报路径越界。")
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                    temporary_path = target_path.with_name(f".{target_path.name}.{uuid.uuid4().hex}.import")
                    shutil.copy2(source_path, temporary_path)
                    temporary_path.replace(target_path)
                    temporary_path = None
                    managed_relative_path = relative.as_posix()
                    external_path = None
                    path_normalized = managed_relative_path.casefold()
                else:
                    managed_relative_path = None
                    external_path = str(source_path)
                    path_normalized = str(source_path).replace("/", "\\").casefold()

                item = LibraryItem(
                    workspace_id=workspace.id,
                    display_name=safe_name,
                    original_name=source_path.name,
                    description=f"{report.title} 第 {version_number} 版",
                    kind="report",
                    storage_mode=storage_mode,
                    managed_relative_path=managed_relative_path,
                    external_path=external_path,
                    path_normalized=path_normalized,
                    mime_type=mime_type,
                    size_bytes=size_bytes,
                    sha256=file_hash,
                    link_status="available",
                    last_verified_at=utcnow(),
                    ai_readability="text" if source_path.suffix.lower() in {".csv", ".md", ".txt"} else "metadata_only",
                )
                db.session.add(item)
                db.session.flush()
                db.session.execute(weekly_report_library_item.insert().values(
                    report_id=report.id,
                    library_item_id=item.id,
                    relation_role=f"version:{version_number:04d}",
                    created_at=utcnow(),
                ))
                if project_id:
                    existing_project_link = db.session.execute(
                        project_library_item.select().where(
                            project_library_item.c.project_id == project_id,
                            project_library_item.c.library_item_id == item.id,
                        )
                    ).first()
                    if not existing_project_link:
                        db.session.execute(project_library_item.insert().values(
                            project_id=project_id,
                            library_item_id=item.id,
                            relation_role="report",
                            created_at=utcnow(),
                        ))

                report.title = _string(payload.get("title") or report.title, 180, "title", required=True)
                report.project_id = project_id
                report.report_date = report_date
                report.period_start = period_start
                report.period_end = period_end
                report.status = _string(payload.get("status") or report.status or "submitted", 20, "status")
                report.summary = _string(payload.get("summary") or report.summary, 200000, "summary")
                report.original_name = source_path.name
                report.stored_path = managed_relative_path or ""
                report.folder_path = (
                    str(Path(managed_relative_path).parent).replace("\\", "/")
                    if managed_relative_path else str(source_path.parent)
                )
                report.mime_type = mime_type
                report.size_bytes = size_bytes
                report.sha256 = file_hash
                self._index_entity(workspace, "file", item)
                self._index_entity(workspace, "weekly", report)
                self._activity(
                    workspace,
                    "weekly_file_imported",
                    f"导入周报文件：{report.title} v{version_number}",
                    "weekly_report",
                    report.id,
                    project_id=project_id,
                )
                db.session.commit()
                return self._weekly_dto(report, detail=True)
            except Exception:
                db.session.rollback()
                if temporary_path:
                    temporary_path.unlink(missing_ok=True)
                if target_path:
                    target_path.unlink(missing_ok=True)
                raise
            finally:
                db.session.remove()

    def _weekly_file_path(self, workspace, payload):
        report_id = _positive_id(payload.get("report_id"), "report_id")
        file_id = _positive_id(payload.get("file_id"), "file_id", optional=True)
        report = WeeklyReport.query.filter_by(
            id=report_id, workspace_id=workspace.id, is_deleted=False,
        ).first()
        if not report:
            raise NotFoundError("周报不存在或已移入回收站。")
        rows = self._weekly_file_rows(report.id)
        if not rows and report.stored_path:
            root = Path(self.flask_app.config["WEEKLY_REPORT_UPLOAD_DIR"]).resolve()
            path = (root / report.stored_path).resolve()
            if root != path and root not in path.parents:
                raise ValidationError("旧周报文件路径越界。")
            if not path.is_file():
                raise NotFoundError("周报文件已移动或不存在。")
            return path, report
        if file_id:
            selected = next(((item, role) for item, role in rows if item.id == file_id), None)
        else:
            selected = max(rows, key=lambda row: self._weekly_version_number(row[1]), default=None)
        if not selected:
            raise NotFoundError("该周报没有可访问的文件。")
        item, _ = selected
        if item.storage_mode == "managed":
            root = (Path(self.flask_app.instance_path) / "library").resolve()
            path = (root / str(item.managed_relative_path or "")).resolve()
            if root != path and root not in path.parents:
                raise ValidationError("托管周报路径越界。")
        else:
            path = Path(item.external_path or "").expanduser().resolve()
        if not path.is_file():
            item.link_status = "missing"
            raise NotFoundError("周报文件已移动或不存在。")
        return path, item

    def weekly_file_path(self, payload):
        def resolve(workspace):
            path, item = self._weekly_file_path(workspace, payload)
            return {"path": str(path), "file_id": item.id, "name": item.original_name}
        return self._run(resolve)

    def weekly_directory_path(self, payload=None):
        payload = payload or {}

        def resolve(workspace):
            report_id = _positive_id(payload.get("report_id"), "report_id", optional=True)
            if report_id:
                path, _ = self._weekly_file_path(workspace, {"report_id": report_id})
                directory = path.parent
            else:
                directory = (Path(self.flask_app.instance_path) / "library" / "weekly").resolve()
                directory.mkdir(parents=True, exist_ok=True)
            return {"path": str(directory)}
        return self._run(resolve)

    def export_weekly_file(self, payload):
        target = Path(_string(payload.get("path"), 2000, "path", required=True)).expanduser().resolve()

        def export(workspace):
            source, item = self._weekly_file_path(workspace, payload)
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(source, target)
            return {"path": str(target), "size_bytes": target.stat().st_size, "name": item.original_name}
        return self._run(export)

    @staticmethod
    def _weekly_dto(item, *, detail=False):
        data = {
            "id": item.id, "project_id": item.project_id, "title": item.title,
            "original_name": item.original_name, "report_date": _iso(item.report_date),
            "period_start": _iso(item.period_start), "period_end": _iso(item.period_end),
            "status": item.status, "summary": item.summary, "row_version": item.row_version,
            "mime_type": item.mime_type, "size_bytes": item.size_bytes,
            "annotation_count": len(item.updates), "updated_at": _iso(item.updated_at),
        }
        if detail:
            files = DesktopModuleServiceMixin._weekly_file_rows(item.id)
            file_items = sorted(
                (DesktopModuleServiceMixin._weekly_file_dto(file_item, role) for file_item, role in files),
                key=lambda value: value["version_number"],
                reverse=True,
            )
            if not file_items and item.stored_path:
                file_items = [{
                    "id": item.id,
                    "version_number": 1,
                    "display_name": item.original_name,
                    "original_name": item.original_name,
                    "storage_mode": "managed",
                    "mime_type": item.mime_type,
                    "size_bytes": item.size_bytes,
                    "sha256": item.sha256,
                    "link_status": "available",
                    "updated_at": _iso(item.updated_at),
                }]
            data.update({
                "body": item.body, "issues_and_feedback": item.issues_and_feedback,
                "next_week_plan": item.next_week_plan, "stored_path": item.stored_path,
                "folder_path": item.folder_path,
                "files": file_items,
                "current_file": file_items[0] if file_items else None,
                "updates": [{"id": update.id, "kind": update.kind, "status": update.status,
                             "content": update.content, "entry_date": _iso(update.entry_date),
                             "created_at": _iso(update.created_at)} for update in item.updates],
                "entries": [{
                    "id": entry.id, "source_type": entry.source_type, "source_id": entry.source_id,
                    "source_title": entry.source_title, "source_excerpt": entry.source_excerpt,
                    "source_date": _iso(entry.source_date), "include_state": entry.include_state,
                } for entry in WeeklyReportEntry.query.filter_by(report_id=item.id).order_by(WeeklyReportEntry.position).all()],
            })
        return data

    def get_weekly(self, payload):
        report_id = _positive_id(payload.get("id"))

        def query(workspace):
            report = WeeklyReport.query.filter_by(id=report_id, workspace_id=workspace.id, is_deleted=False).first()
            if not report:
                raise NotFoundError("周报不存在或已移入回收站。")
            return self._weekly_dto(report, detail=True)
        return self._run(query)

    def add_weekly_annotation(self, payload, expected_row_version=None):
        report_id = _positive_id(payload.get("id"))
        content = _string(payload.get("content"), 100000, "content", required=True)
        kind = _string(payload.get("kind") or "批注", 20, "kind")
        status = _string(payload.get("status") or "待处理", 20, "status")
        if kind not in {"批注", "说明", "指导"}:
            raise ValidationError("周报批注类型无效。")
        if status not in {"待处理", "已处理"}:
            raise ValidationError("周报批注状态无效。")

        def save(workspace):
            report = WeeklyReport.query.filter_by(
                id=report_id, workspace_id=workspace.id, is_deleted=False,
            ).first()
            if not report:
                raise NotFoundError("周报不存在或已移入回收站。")
            if expected_row_version is None:
                raise ValidationError("缺少周报版本。")
            if report.row_version != int(expected_row_version):
                raise ConflictError("周报已被更新，请重新加载。")
            db.session.add(WeeklyReportUpdate(
                report_id=report.id, user_id=workspace.legacy_user_id,
                entry_date=date.today(), kind=kind, status=status, content=content,
            ))
            report.row_version += 1
            db.session.flush()
            return self._weekly_dto(report, detail=True)
        return self._run(save)

    def update_weekly(self, payload, expected_row_version=None):
        report_id = _positive_id(payload.get("id"))

        def update_report(workspace):
            report = WeeklyReport.query.filter_by(id=report_id, workspace_id=workspace.id, is_deleted=False).first()
            if not report:
                raise NotFoundError("周报不存在或已移入回收站。")
            if expected_row_version is None:
                raise ValidationError("缺少周报版本。")
            if report.row_version != int(expected_row_version):
                raise ConflictError("周报已被更新，请重新加载。")
            if "project_id" in payload:
                project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
                self._ensure_project(workspace.id, project_id)
                report.project_id = project_id
            for field, maximum in (("title", 180), ("status", 20), ("summary", 200000)):
                if field in payload:
                    setattr(report, field, _string(payload.get(field), maximum, field, required=field == "title"))
            report.row_version += 1
            db.session.flush()
            self._index_entity(workspace, "weekly", report)
            return self._weekly_dto(report, detail=True)
        return self._run(update_report)

    def list_weekly(self, payload=None):
        payload = payload or {}
        search = _string(payload.get("search"), 180, "search")
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        status = _string(payload.get("status"), 20, "status")
        date_start = _date(payload.get("date_start"), "date_start")
        date_end = _date(payload.get("date_end"), "date_end")
        if date_start and date_end and date_end < date_start:
            raise ValidationError("结束日期不能早于开始日期。")

        def query(workspace):
            reports = WeeklyReport.query.filter_by(workspace_id=workspace.id, is_deleted=False)
            if project_id:
                self._ensure_project(workspace.id, project_id)
                reports = reports.filter(WeeklyReport.project_id == project_id)
            if status:
                reports = reports.filter(WeeklyReport.status == status)
            if date_start:
                reports = reports.filter(WeeklyReport.report_date >= date_start)
            if date_end:
                reports = reports.filter(WeeklyReport.report_date <= date_end)
            if search:
                escaped = search.replace("%", r"\%").replace("_", r"\_")
                reports = reports.filter(WeeklyReport.title.ilike(f"%{escaped}%", escape="\\"))
            reports = _sorted(reports, payload, {
                "date_desc": (WeeklyReport.report_date.desc(), WeeklyReport.updated_at.desc()),
                "updated_desc": (WeeklyReport.updated_at.desc(),), "title_asc": (WeeklyReport.title.asc(),),
            }, "date_desc")
            return _page_result(reports, payload, self._weekly_dto)
        return self._run(query)

    def weekly_bulk(self, payload):
        ids = sorted({_positive_id(value) for value in (payload.get("ids") or [])})
        if not ids or len(ids) > 500:
            raise ValidationError("请选择 1 至 500 份周报。")
        action = _string(payload.get("action"), 20, "action", required=True)
        if action not in {"status", "project", "trash"}:
            raise ValidationError("不支持的周报批量操作。")

        def bulk(workspace):
            items = WeeklyReport.query.filter(WeeklyReport.workspace_id == workspace.id, WeeklyReport.is_deleted.is_(False), WeeklyReport.id.in_(ids)).all()
            if action == "status":
                value = _string(payload.get("value"), 20, "value", required=True)
                if value not in {"draft", "submitted", "reviewed", "archived"}: raise ValidationError("周报状态无效。")
            elif action == "project":
                value = _positive_id(payload.get("value"), "value", optional=True)
                self._ensure_project(workspace.id, value)
            else: value = None
            results = []
            for item in items:
                if action == "trash": item.is_deleted, item.deleted_at = True, utcnow()
                elif action == "status": item.status = value
                else: item.project_id = value
                item.row_version += 1
                self._index_entity(workspace, "weekly", item)
                results.append({"id": item.id, "status": "updated"})
            found = {item.id for item in items}
            results.extend({"id": value, "status": "not_found"} for value in ids if value not in found)
            return {"updated": len(items), "skipped": len(ids) - len(items), "results": results}
        return self._run(bulk)

    def ai_preview(self, payload):
        target_type = _string(payload.get("target_type"), 30, "target_type", required=True)
        target_id = _positive_id(payload.get("target_id"), "target_id")
        source_ids = payload.get("source_ids") if isinstance(payload.get("source_ids"), list) else []
        records_explicit = bool(payload.get("records_explicit"))

        def preview(workspace):
            preset = ApiPreset.query.filter_by(user_id=workspace.legacy_user_id, is_default=True).first()
            if not preset:
                preset = ApiPreset.query.filter_by(user_id=workspace.legacy_user_id).first()
            target = self._ai_target(workspace, target_type, target_id)
            target_snapshot = self._ai_target_snapshot(target_type, target)
            project_records = []
            if target_type == "project":
                # Project chat deliberately reads only structured experiment records, never library files.
                records = LabRecord.query.filter_by(
                    workspace_id=workspace.id, project_id=target.id, is_deleted=False,
                )
                if source_ids:
                    selected_ids = list(dict.fromkeys(_positive_id(value, "source_ids") for value in source_ids))[:200]
                    records = records.filter(LabRecord.id.in_(selected_ids))
                elif records_explicit:
                    records = records.filter(False)
                project_records = [self._record_dto(record, detail=True) for record in records.order_by(
                    LabRecord.experiment_date.desc(), LabRecord.updated_at.desc(),
                ).all()]
                if source_ids and len(project_records) != len(selected_ids):
                    raise ValidationError("部分实验记录不属于当前项目。")
            source_snapshot = {
                "target_type": target_type, "target_id": target_id, "source_ids": source_ids,
                "target_snapshot": target_snapshot, "project_records": project_records,
                "content_hash": hashlib.sha256(_json({
                    "target": target_snapshot, "records": project_records, "sources": source_ids,
                }).encode("utf-8")).hexdigest(),
            }
            return {
                "target": target_snapshot,
                "base_row_version": target.row_version,
                "endpoint": preset.api_url if preset else "未配置",
                "model": preset.text_model if preset else "未配置",
                "api_enabled": bool(preset and preset.is_enabled and preset.encrypted_api_key),
                "source_snapshot": source_snapshot,
                "estimated_characters": len(_json(source_snapshot)),
                "warning": "应用操作会直接写入所选字段，并生成可撤销的变更记录。",
            }
        return self._run(preview)

    def ai_history(self, payload):
        try:
            page = max(1, int(payload.get("page") or 1))
            per_page = int(payload.get("per_page") or 5)
        except (TypeError, ValueError) as exc:
            raise ValidationError("提示词历史分页参数无效。") from exc
        if per_page not in {5, 10, 20}:
            raise ValidationError("提示词历史每页条数无效。")
        query_text = _string(payload.get("query"), 200, "query")

        def query(workspace):
            history = (
                db.session.query(AIMessage, AIConversation, AIChangeSet)
                .join(AIConversation, AIMessage.conversation_id == AIConversation.id)
                .join(AIChangeSet, AIChangeSet.message_id == AIMessage.id)
                .filter(
                    AIConversation.user_id == workspace.legacy_user_id,
                    AIMessage.role == "assistant",
                    AIMessage.prompt_snapshot != "",
                )
            )
            if query_text:
                escaped = query_text.replace("%", r"\%").replace("_", r"\_")
                history = history.filter(AIMessage.prompt_snapshot.ilike(f"%{escaped}%", escape="\\"))
            total = history.count()
            pages = max(1, (total + per_page - 1) // per_page)
            page_number = min(page, pages)
            rows = (
                history.order_by(AIMessage.created_at.desc(), AIMessage.id.desc())
                .offset((page_number - 1) * per_page)
                .limit(per_page)
                .all()
            )
            return {
                "items": [{
                    "id": message.id,
                    "prompt": message.prompt_snapshot,
                    "reply": message.content,
                    "target_type": change.target_type if change else conversation.page_type,
                    "target_id": change.target_id if change else conversation.page_id,
                    "status": change.status if change else "generated",
                    "model": message.model_name,
                    "created_at": _iso(message.created_at),
                } for message, conversation, change in rows],
                "pagination": {
                    "page": page_number, "pages": pages, "per_page": per_page, "total": total,
                },
            }
        return self._run(query)

    @staticmethod
    def _ai_conversation_or_error(workspace, conversation_id):
        conversation = AIConversation.query.filter_by(
            id=conversation_id, user_id=workspace.legacy_user_id,
        ).first()
        if not conversation:
            raise NotFoundError("聊天记录不存在。")
        return conversation

    def ai_conversations(self, payload=None):
        payload = payload or {}
        try:
            page = max(1, int(payload.get("page") or 1))
            per_page = int(payload.get("per_page") or 10)
        except (TypeError, ValueError) as exc:
            raise ValidationError("聊天记录分页参数无效。") from exc
        if per_page not in {5, 10, 20}:
            raise ValidationError("聊天记录每页条数无效。")
        query_text = _string(payload.get("query"), 200, "query")
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        date_start = _date(payload.get("date_start"), "date_start")
        date_end = _date(payload.get("date_end"), "date_end")
        if date_start and date_end and date_end < date_start:
            raise ValidationError("聊天结束日期不能早于开始日期。")

        def query(workspace):
            conversations = AIConversation.query.filter_by(user_id=workspace.legacy_user_id)
            if project_id:
                self._ensure_project(workspace.id, project_id)
                conversations = conversations.filter_by(page_type="project", page_id=project_id)
            if date_start:
                conversations = conversations.filter(AIConversation.created_at >= datetime.combine(date_start, time.min))
            if date_end:
                conversations = conversations.filter(AIConversation.created_at <= datetime.combine(date_end, time.max))
            if query_text:
                escaped = query_text.replace("%", r"\%").replace("_", r"\_")
                matching_ids = db.select(AIMessage.conversation_id).where(or_(
                    AIMessage.content.ilike(f"%{escaped}%", escape="\\"),
                    AIMessage.prompt_snapshot.ilike(f"%{escaped}%", escape="\\"),
                ))
                conversations = conversations.filter(or_(
                    AIConversation.title.ilike(f"%{escaped}%", escape="\\"),
                    AIConversation.id.in_(matching_ids),
                ))
            total = conversations.count()
            pages = max(1, (total + per_page - 1) // per_page)
            page_number = min(page, pages)
            rows = conversations.order_by(AIConversation.updated_at.desc(), AIConversation.id.desc()).offset(
                (page_number - 1) * per_page
            ).limit(per_page).all()
            items = []
            for conversation in rows:
                messages = list(conversation.messages)
                assistant_message = next((message for message in reversed(messages) if message.role == "assistant"), None)
                items.append({
                    "id": conversation.id, "title": conversation.title, "project_id": conversation.page_id
                    if conversation.page_type == "project" else None,
                    "message_count": len(messages), "preview": assistant_message.content if assistant_message else "",
                    "updated_at": _iso(conversation.updated_at), "created_at": _iso(conversation.created_at),
                })
            return {"items": items, "pagination": {
                "page": page_number, "pages": pages, "per_page": per_page, "total": total,
            }}
        return self._run(query)

    def ai_conversation_get(self, payload):
        conversation_id = _positive_id(payload.get("id"))

        def query(workspace):
            conversation = self._ai_conversation_or_error(workspace, conversation_id)
            message_ids = [message.id for message in conversation.messages]
            changes = {change.message_id: change for change in AIChangeSet.query.filter(
                AIChangeSet.message_id.in_(message_ids) if message_ids else False
            ).all()} if message_ids else {}
            return {
                "id": conversation.id, "title": conversation.title, "project_id": conversation.page_id
                if conversation.page_type == "project" else None,
                "selected_record_ids": _load_json(conversation.selected_record_ids_json, []),
                "messages": [{
                    "id": message.id, "role": message.role, "content": message.content,
                    "prompt": message.prompt_snapshot, "model": message.model_name,
                    "created_at": _iso(message.created_at),
                    "change_id": changes[message.id].id if message.id in changes else None,
                    "change_status": changes[message.id].status if message.id in changes else None,
                    "proposal": _load_json(message.proposal_json, {}),
                } for message in conversation.messages],
            }
        return self._run(query)

    def ai_conversation_create(self, payload=None):
        payload = payload or {}
        project_id = _positive_id(payload.get("project_id"), "project_id", optional=True)
        title = _string(payload.get("title") or "新聊天", 160, "title")
        record_ids = payload.get("record_ids") if isinstance(payload.get("record_ids"), list) else []
        selected_record_ids = list(dict.fromkeys(_positive_id(value, "record_ids") for value in record_ids))[:200]

        def create(workspace):
            if project_id:
                self._ensure_project(workspace.id, project_id)
                if selected_record_ids:
                    count = LabRecord.query.filter(
                        LabRecord.workspace_id == workspace.id, LabRecord.project_id == project_id,
                        LabRecord.is_deleted.is_(False), LabRecord.id.in_(selected_record_ids),
                    ).count()
                    if count != len(selected_record_ids):
                        raise ValidationError("部分实验记录不属于当前项目。")
            conversation = AIConversation(
                user_id=workspace.legacy_user_id, title=title,
                page_type="project" if project_id else "", page_id=project_id,
                selected_record_ids_json=_json(selected_record_ids),
            )
            db.session.add(conversation)
            db.session.flush()
            return {"id": conversation.id, "title": conversation.title, "project_id": project_id}
        return self._run(create)

    def ai_conversations_bulk(self, payload):
        ids = payload.get("ids") if isinstance(payload.get("ids"), list) else []
        conversation_ids = list(dict.fromkeys(_positive_id(value, "ids") for value in ids))[:100]
        if not conversation_ids:
            raise ValidationError("请至少选择一条聊天记录。")
        action = _string(payload.get("action"), 20, "action", required=True)
        title = _string(payload.get("title"), 160, "title")
        if action not in {"rename", "delete"}:
            raise ValidationError("不支持的聊天记录批量操作。")
        if action == "rename" and not title:
            raise ValidationError("请输入新的聊天标题。")

        def bulk(workspace):
            conversations = AIConversation.query.filter(
                AIConversation.user_id == workspace.legacy_user_id,
                AIConversation.id.in_(conversation_ids),
            ).all()
            if len(conversations) != len(conversation_ids):
                raise NotFoundError("部分聊天记录不存在。")
            if action == "rename":
                for index, conversation in enumerate(conversations, start=1):
                    conversation.title = title if len(conversations) == 1 else f"{title} {index}"
                return {"updated": len(conversations), "skipped": 0}
            deleted = skipped = 0
            for conversation in conversations:
                message_ids = [message.id for message in conversation.messages]
                changes = AIChangeSet.query.filter(AIChangeSet.message_id.in_(message_ids)).all() if message_ids else []
                if any(change.status in {"applied", "partially_applied"} for change in changes):
                    skipped += 1
                    continue
                for change in changes:
                    db.session.delete(change)
                db.session.delete(conversation)
                deleted += 1
            return {"updated": deleted, "skipped": skipped}
        return self._run(bulk)

    @staticmethod
    def _ai_target(workspace, target_type, target_id):
        models = {
            "project": ResearchProject, "lab_record": LabRecord, "task": Task,
            "calendar_event": CalendarEvent, "note": Note, "weekly_report": WeeklyReport,
        }
        model = models.get(target_type)
        if not model:
            raise ValidationError("AI 目标类型无效。")
        target = model.query.filter_by(id=target_id, workspace_id=workspace.id, is_deleted=False).first()
        if not target:
            raise NotFoundError("AI 写入目标不存在。")
        return target

    @staticmethod
    def _ai_target_snapshot(target_type, target):
        snapshot = {}
        for field in AI_TARGET_FIELDS[target_type]:
            if field == "steps":
                value = [{
                    "title": step.title,
                    "instruction": step.instruction,
                    "planned_duration_minutes": step.planned_duration_minutes,
                    "checkpoint": step.checkpoint,
                    "risk": step.risk,
                    "is_done": step.is_done,
                    "actual_deviation": step.actual_deviation,
                } for step in target.steps]
            elif field == "annotation":
                latest = WeeklyReportUpdate.query.filter_by(report_id=target.id, kind="AI批注").order_by(
                    WeeklyReportUpdate.created_at.desc(), WeeklyReportUpdate.id.desc()
                ).first()
                value = latest.content if latest else ""
            else:
                value = getattr(target, field)
            snapshot[field] = _iso(value) if isinstance(value, (date, datetime)) else (value if value is not None else "")
        return snapshot

    @staticmethod
    def _ai_field_schema(target_type):
        schema = {}
        for field in sorted(AI_TARGET_FIELDS[target_type]):
            field_type = "text"
            if field in AI_DATE_FIELDS:
                field_type = "date:YYYY-MM-DD"
            elif field == "steps":
                field_type = (
                    "array:{title:string,instruction:string,planned_duration_minutes:integer-or-null,"
                    "checkpoint:string,risk:string,is_done:boolean,actual_deviation:string}"
                )
            elif field in AI_DATETIME_FIELDS:
                field_type = "datetime:ISO-8601"
            elif field in AI_ID_FIELDS:
                field_type = "positive-integer-or-null"
            elif field in AI_BOOLEAN_FIELDS:
                field_type = "boolean"
            elif (target_type, field) in AI_ENUM_FIELDS:
                field_type = "enum:" + "|".join(sorted(AI_ENUM_FIELDS[(target_type, field)]))
            schema[field] = {"label": AI_FIELD_LABELS.get(field, field), "type": field_type}
        return schema

    @staticmethod
    def _normalize_ai_value(target_type, field, value):
        if field == "steps":
            if target_type != "lab_record" or not isinstance(value, list) or len(value) > 200:
                raise ValidationError("AI 返回的实验步骤无效。", field_errors={field: "必须为不超过 200 项的数组"})
            normalized_steps = []
            for index, raw in enumerate(value, start=1):
                if not isinstance(raw, dict):
                    raise ValidationError("AI 返回的实验步骤无效。", field_errors={field: f"第 {index} 项必须为对象"})
                duration = raw.get("planned_duration_minutes")
                if duration in (None, ""):
                    duration = None
                else:
                    try:
                        duration = int(duration)
                    except (TypeError, ValueError) as exc:
                        raise ValidationError(
                            "AI 返回的实验步骤无效。",
                            field_errors={field: f"第 {index} 项预计耗时必须为整数"},
                        ) from exc
                    if duration < 0:
                        raise ValidationError(
                            "AI 返回的实验步骤无效。",
                            field_errors={field: f"第 {index} 项预计耗时不能小于 0"},
                        )
                done = raw.get("is_done", False)
                if not isinstance(done, bool):
                    raise ValidationError(
                        "AI 返回的实验步骤无效。", field_errors={field: f"第 {index} 项完成状态必须为布尔值"},
                    )
                normalized_steps.append({
                    "title": _string(raw.get("title"), 180, f"steps[{index}].title", required=True),
                    "instruction": _string(raw.get("instruction"), 20000, f"steps[{index}].instruction"),
                    "planned_duration_minutes": duration,
                    "checkpoint": _string(raw.get("checkpoint"), 5000, f"steps[{index}].checkpoint"),
                    "risk": _string(raw.get("risk"), 5000, f"steps[{index}].risk"),
                    "is_done": done,
                    "actual_deviation": _string(raw.get("actual_deviation"), 5000, f"steps[{index}].actual_deviation"),
                })
            return normalized_steps
        if field in AI_DATE_FIELDS:
            return _iso(_date(value, field)) if value not in (None, "") else ""
        if field in AI_DATETIME_FIELDS:
            return _iso(_datetime(value, field)) if value not in (None, "") else ""
        if field in AI_ID_FIELDS:
            return _positive_id(value, field, optional=True)
        if field in AI_BOOLEAN_FIELDS:
            if isinstance(value, bool):
                return value
            normalized = str(value or "").strip().lower()
            if normalized in {"true", "1", "yes", "是", "全天"}:
                return True
            if normalized in {"false", "0", "no", "否", ""}:
                return False
            raise ValidationError("AI 返回的布尔字段无效。", field_errors={field: "应为 true 或 false"})
        maximum = AI_FIELD_MAX_LENGTH.get(field, 200000)
        normalized = _string(value, maximum, field, required=field == "title")
        choices = AI_ENUM_FIELDS.get((target_type, field))
        if choices and normalized not in choices:
            raise ValidationError("AI 返回的枚举字段无效。", field_errors={field: "不支持的选项"})
        return normalized

    @staticmethod
    def _ai_assignment_value(field, value):
        if field in AI_DATE_FIELDS:
            return _date(value, field)
        if field in AI_DATETIME_FIELDS:
            return _datetime(value, field) if value else None
        return value

    def _validate_ai_relations(self, workspace, target, values):
        project_id = values.get("project_id", getattr(target, "project_id", None))
        record_id = values.get("lab_record_id", getattr(target, "lab_record_id", None))
        self._ensure_project(workspace.id, project_id)
        if record_id:
            record = LabRecord.query.filter_by(id=record_id, workspace_id=workspace.id, is_deleted=False).first()
            if not record:
                raise NotFoundError("AI 关联的实验记录不存在。")
            if project_id and record.project_id != project_id:
                raise ValidationError("AI 关联的项目与实验记录不一致。")

    def _apply_ai_field(
        self, workspace, target_type, target, field, value, *, source_ai_message_id=None, source_kind="ai",
    ):
        if target_type == "lab_record" and field == "steps":
            LabRecordStep.query.filter_by(record_id=target.id).delete(synchronize_session=False)
            db.session.flush()
            db.session.expire(target, ["steps"])
            for position, step in enumerate(value, start=1):
                target.steps.append(LabRecordStep(
                    position=position, title=step["title"], instruction=step["instruction"],
                    planned_duration_minutes=step["planned_duration_minutes"], checkpoint=step["checkpoint"],
                    risk=step["risk"], is_done=step["is_done"], actual_deviation=step["actual_deviation"],
                    source_kind=source_kind, source_ai_message_id=source_ai_message_id,
                ))
            return None
        if target_type == "weekly_report" and field == "annotation":
            update_item = WeeklyReportUpdate(
                report_id=target.id, user_id=workspace.legacy_user_id, entry_date=date.today(),
                kind="AI批注", status="待处理", content=value,
            )
            db.session.add(update_item)
            db.session.flush()
            return update_item.id
        setattr(target, field, self._ai_assignment_value(field, value))
        return None

    def ai_propose(self, payload):
        target_type = _string(payload.get("target_type"), 30, "target_type", required=True)
        target_id = _positive_id(payload.get("target_id"), "target_id")
        prompt = _string(payload.get("prompt"), 30000, "prompt", required=True)
        source_snapshot = payload.get("source_snapshot") if isinstance(payload.get("source_snapshot"), dict) else {}
        conversation_id = _positive_id(payload.get("conversation_id"), "conversation_id", optional=True)
        web_access = bool(payload.get("web_access"))
        record_ids = payload.get("record_ids") if isinstance(payload.get("record_ids"), list) else []
        selected_record_ids = list(dict.fromkeys(_positive_id(value, "record_ids") for value in record_ids))[:200]

        def propose(workspace):
            target = self._ai_target(workspace, target_type, target_id)
            preset = ApiPreset.query.filter_by(user_id=workspace.legacy_user_id, is_default=True).first()
            if not preset:
                preset = ApiPreset.query.filter_by(user_id=workspace.legacy_user_id).first()
            if not preset or not preset.is_enabled or not preset.encrypted_api_key:
                raise ValidationError("请先在设置中启用 API 并配置 Key。")
            from app.ai_service import AIConfig, AIServiceError, chat_with_assistant
            target_snapshot = self._ai_target_snapshot(target_type, target)
            field_schema = self._ai_field_schema(target_type)
            allowed = sorted(AI_TARGET_FIELDS[target_type])
            system_prompt = build_desktop_ai_system_prompt(
                target_type, field_schema, today=date.today(), timezone=workspace.timezone,
            )
            config = AIConfig(
                api_url=preset.api_url, api_key=preset.get_api_key(), model=preset.text_model,
                enabled=preset.is_enabled, allow_private=bool(self.flask_app.config.get("ALLOW_PRIVATE_API_URLS")),
                allowed_hosts=tuple(self.flask_app.config.get("AI_ALLOWED_HOSTS") or ()), source="desktop",
            )
            try:
                user_content = _json({
                    "request": prompt, "current_target": target_snapshot,
                    "source_snapshot": source_snapshot,
                })
                kwargs = {"config": config}
                if web_access:
                    kwargs["web_access"] = True
                result = chat_with_assistant([{"role": "user", "content": user_content}], system_prompt, **kwargs)
            except AIServiceError as exc:
                raise ValidationError(str(exc)) from exc
            proposal = result.get("proposal") if isinstance(result.get("proposal"), dict) else {}
            proposal = {
                field: self._normalize_ai_value(target_type, field, value)
                for field, value in proposal.items() if field in allowed
            }
            if conversation_id:
                conversation = self._ai_conversation_or_error(workspace, conversation_id)
                if conversation.page_type not in {"", target_type} or conversation.page_id not in {None, target_id}:
                    raise ValidationError("该聊天记录不属于当前项目。")
                conversation.page_type, conversation.page_id, conversation.updated_at = target_type, target_id, utcnow()
                conversation.selected_record_ids_json = _json(selected_record_ids)
                if conversation.title in {"", "新聊天"}:
                    conversation.title = prompt[:80]
            else:
                conversation = AIConversation(
                    user_id=workspace.legacy_user_id, title=prompt[:80], page_type=target_type, page_id=target_id,
                    selected_record_ids_json=_json(selected_record_ids),
                )
                db.session.add(conversation)
                db.session.flush()
            db.session.add(AIMessage(
                conversation_id=conversation.id, role="user", content=prompt,
                prompt_snapshot="", context_snapshot_json=_json(source_snapshot),
            ))
            message = AIMessage(
                conversation_id=conversation.id, role="assistant", content=result.get("reply") or "已生成字段建议。",
                proposal_json=_json(proposal), model_name=preset.text_model, prompt_snapshot=prompt,
                context_snapshot_json=_json(source_snapshot), requires_human_review=False,
            )
            db.session.add(message)
            db.session.flush()
            change = AIChangeSet(
                message_id=message.id, target_type=target_type, target_id=target_id,
                base_row_version=target.row_version, proposal_json=_json(proposal),
                before_json=_json(self._ai_target_snapshot(target_type, target)),
                source_snapshot_json=_json(source_snapshot), model_name=preset.text_model, prompt_snapshot=prompt,
            )
            db.session.add(change)
            db.session.flush()
            return {"id": change.id, "conversation_id": conversation.id, "reply": message.content,
                    "proposal": proposal, "field_schema": field_schema, "status": change.status,
                    "web_used": bool(result.get("web_used")), "references": result.get("references") or []}
        return self._run(propose)

    def ai_apply(self, payload):
        change_id = _positive_id(payload.get("id"))
        accepted_input = payload.get("accepted_fields") if isinstance(payload.get("accepted_fields"), list) else []

        def apply(workspace):
            change = AIChangeSet.query.filter_by(id=change_id).first()
            if not change or change.status not in {"proposed", "stale"}:
                raise NotFoundError("AI 变更集不存在或已处理。")
            target = self._ai_target(workspace, change.target_type, change.target_id)
            if target.row_version != change.base_row_version:
                change.status = "stale"
                raise ConflictError("目标已发生变化，请重新生成差异。")
            proposal = _load_json(change.proposal_json, {})
            accepted = [field for field in accepted_input if field in proposal and field in AI_TARGET_FIELDS[change.target_type]]
            if not accepted:
                raise ValidationError("请至少选择一个需要写入的字段。")
            normalized = {
                field: self._normalize_ai_value(change.target_type, field, proposal[field]) for field in accepted
            }
            self._validate_ai_relations(workspace, target, normalized)
            before = self._ai_target_snapshot(change.target_type, target)
            created_updates = {}
            for field, value in normalized.items():
                created_id = self._apply_ai_field(
                    workspace, change.target_type, target, field, value,
                    source_ai_message_id=change.message_id,
                )
                if created_id:
                    created_updates[field] = created_id
            target.row_version += 1
            after = self._ai_target_snapshot(change.target_type, target)
            if created_updates:
                after["_created_update_ids"] = created_updates
            change.before_json, change.after_json = _json(before), _json(after)
            change.accepted_fields_json = _json(accepted)
            change.rejected_fields_json = _json(sorted(set(proposal) - set(accepted)))
            change.status = "applied" if len(accepted) == len(proposal) else "partially_applied"
            change.reviewed_at = change.applied_at = utcnow()
            if change.target_type == "lab_record":
                db.session.add(LabRecordRevision(
                    record_id=target.id, scope="ai_change", reason="应用 AI 字段建议",
                    before_json=_json(before), after_json=_json(after),
                    diff_json=_json({field: proposal[field] for field in accepted}),
                    actor_kind="ai_assisted", source_ai_change_set_id=change.id,
                ))
            self._index_entity(workspace, "weekly" if change.target_type == "weekly_report" else change.target_type.replace("lab_", ""), target)
            self._activity(
                workspace, "ai_change_applied", f"AI 写入：{getattr(target, 'title', change.target_type)}",
                target.__tablename__, target.id, project_id=getattr(target, "project_id", None),
                record_id=target.id if change.target_type == "lab_record" else getattr(target, "lab_record_id", None),
            )
            return {"id": change.id, "status": change.status, "row_version": target.row_version, "applied_fields": accepted}
        return self._run(apply)

    def ai_revert(self, payload):
        change_id = _positive_id(payload.get("id"))

        def revert(workspace):
            change = AIChangeSet.query.filter_by(id=change_id).first()
            if not change or change.status not in {"applied", "partially_applied"}:
                raise NotFoundError("该 AI 变更不可撤销。")
            target = self._ai_target(workspace, change.target_type, change.target_id)
            if target.row_version != change.base_row_version + 1:
                raise ConflictError("目标在 AI 应用后又发生了变化，不能直接撤销，请先重新加载。")
            before_apply = _load_json(change.before_json, {})
            after_apply = _load_json(change.after_json, {})
            current = self._ai_target_snapshot(change.target_type, target)
            for field in _load_json(change.accepted_fields_json, []):
                if field == "annotation":
                    update_id = (after_apply.get("_created_update_ids") or {}).get(field)
                    if update_id:
                        update_item = WeeklyReportUpdate.query.filter_by(id=update_id, report_id=target.id).first()
                        if update_item:
                            db.session.delete(update_item)
                elif field in before_apply:
                    self._apply_ai_field(
                        workspace, change.target_type, target, field, before_apply[field], source_kind="local_user",
                    )
            target.row_version += 1
            restored = self._ai_target_snapshot(change.target_type, target)
            change.status, change.reverted_at = "reverted", utcnow()
            if change.target_type == "lab_record":
                db.session.add(LabRecordRevision(
                    record_id=target.id, scope="ai_revert", reason="撤销 AI 字段建议",
                    before_json=_json(current), after_json=_json(restored), diff_json=change.before_json,
                    actor_kind="local_user", source_ai_change_set_id=change.id,
                ))
            self._index_entity(workspace, "weekly" if change.target_type == "weekly_report" else change.target_type.replace("lab_", ""), target)
            return {"id": change.id, "status": change.status, "row_version": target.row_version}
        return self._run(revert)

    def rebuild_search(self, payload=None):
        def rebuild(workspace):
            SearchDocument.query.filter_by(workspace_id=workspace.id).delete()
            documents = []
            for item in ResearchProject.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                documents.append(("project", item.id, item.title, f"{item.code or ''}\n{item.objective or ''}\n{item.notes or ''}", "projects"))
            for item in LabRecord.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                body = "\n".join(str(getattr(item, field) or "") for field in (
                    "record_code", "objective", "background", "hypothesis", "design", "materials_conditions",
                    "expected_result", "actual_process_summary", "actual_result", "analysis", "conclusion", "next_steps",
                ))
                documents.append(("record", item.id, item.title, body, "record-edit"))
            for item in LiteratureItem.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                documents.append(("literature", item.id, item.title, f"{item.authors_json}\n{item.doi}\n{item.abstract}\n{item.reading_notes}", "literature"))
            for item in LibraryItem.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                documents.append(("file", item.id, item.display_name, f"{item.external_path or item.managed_relative_path or ''}\n{item.description}", "files"))
            for item in Note.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                documents.append(("note", item.id, item.title, item.body, "notes"))
            for item in Task.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                documents.append(("task", item.id, item.title, item.notes or "", "tasks"))
            for item in WeeklyReport.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                documents.append(("weekly", item.id, item.title, f"{item.body}\n{item.issues_and_feedback}\n{item.next_week_plan}", "weekly"))
            for entity_type, entity_id, title_value, body, view_key in documents:
                db.session.add(SearchDocument(
                    workspace_id=workspace.id, entity_type=entity_type, entity_id=entity_id,
                    title=title_value, body=body, keywords="", view_key=view_key,
                ))
            db.session.flush()
            mode = "like"
            try:
                db.session.execute(text("INSERT INTO search_fts(search_fts) VALUES ('rebuild')"))
                mode = "fts5-trigram"
            except Exception:
                db.session.rollback()
                # The derived documents are rebuilt after rollback so LIKE fallback remains complete.
                SearchDocument.query.filter_by(workspace_id=workspace.id).delete()
                for entity_type, entity_id, title_value, body, view_key in documents:
                    db.session.add(SearchDocument(
                        workspace_id=workspace.id, entity_type=entity_type, entity_id=entity_id,
                        title=title_value, body=body, keywords="", view_key=view_key,
                    ))
                db.session.flush()
            return {"count": len(documents), "mode": mode,
                    "diagnostic": "FTS5 trigram" if mode == "fts5-trigram" else "FTS5 trigram 不可用，当前使用 LIKE 降级搜索"}
        return self._run(rebuild)

    def search(self, payload=None):
        payload = payload or {}
        query_text = _string(payload.get("query"), 200, "query")
        if not query_text:
            return {"items": [], "mode": "idle", "diagnostic": "输入关键词开始搜索"}

        def query(workspace):
            if not SearchDocument.query.filter_by(workspace_id=workspace.id).first():
                # Inline rebuild avoids a second application-context transaction.
                for item in ResearchProject.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                    db.session.add(SearchDocument(workspace_id=workspace.id, entity_type="project", entity_id=item.id,
                                                  title=item.title, body=item.objective or "", keywords=item.code or "", view_key="projects"))
                for item in LabRecord.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                    body = "\n".join((item.objective, item.design, item.actual_result, item.conclusion, item.next_steps))
                    db.session.add(SearchDocument(workspace_id=workspace.id, entity_type="record", entity_id=item.id,
                                                  title=item.title, body=body, keywords=item.record_code, view_key="record-edit"))
                for model, entity_type, view_key, body_field in (
                    (LiteratureItem, "literature", "literature", "abstract"),
                    (LibraryItem, "file", "files", "description"),
                    (Note, "note", "notes", "body"),
                    (Task, "task", "tasks", "notes"),
                    (WeeklyReport, "weekly", "weekly", "body"),
                ):
                    for item in model.query.filter_by(workspace_id=workspace.id, is_deleted=False):
                        title_value = getattr(item, "title", None) or getattr(item, "display_name", "")
                        db.session.add(SearchDocument(workspace_id=workspace.id, entity_type=entity_type, entity_id=item.id,
                                                      title=title_value, body=getattr(item, body_field) or "",
                                                      keywords="", view_key=view_key))
                db.session.flush()
            mode, diagnostic = "like", "FTS5 trigram 不可用，当前使用 LIKE 降级搜索"
            rows = []
            if len(query_text) >= 3:
                try:
                    match_query = f'"{query_text.replace(chr(34), chr(34) * 2)}"'
                    rows = db.session.execute(text(
                        "SELECT sd.* FROM search_fts "
                        "JOIN search_document sd ON sd.id = search_fts.rowid "
                        "WHERE search_fts MATCH :query AND sd.workspace_id = :workspace_id "
                        "ORDER BY rank LIMIT 80"
                    ), {"query": match_query, "workspace_id": workspace.id}).mappings().all()
                    mode, diagnostic = "fts5-trigram", "FTS5 trigram 连续子串匹配"
                except Exception:
                    rows = []
            if mode == "like":
                pattern = f"%{query_text}%"
                rows = SearchDocument.query.filter(
                    SearchDocument.workspace_id == workspace.id,
                    or_(SearchDocument.title.ilike(pattern), SearchDocument.body.ilike(pattern), SearchDocument.keywords.ilike(pattern)),
                ).order_by(SearchDocument.updated_at.desc()).limit(80).all()
            return {"mode": mode, "diagnostic": diagnostic,
                    "items": [{"id": row.entity_id, "entity_type": row.entity_type, "title": row.title,
                               "excerpt": re.sub(r"\s+", " ", row.body)[:180], "view_key": row.view_key}
                              for row in rows]}
        return self._run(query)

    def trash_list(self, payload=None):
        def query(workspace):
            items = []
            for model, entity_type, label_field in (
                (ResearchProject, "project", "title"), (LabRecord, "lab_record", "title"),
                (LiteratureItem, "literature", "title"), (LibraryItem, "library_item", "display_name"),
                (Note, "note", "title"), (Task, "task", "title"), (WeeklyReport, "weekly_report", "title"),
                (CalendarEvent, "calendar_event", "title"),
            ):
                for item in model.query.filter_by(workspace_id=workspace.id, is_deleted=True).order_by(model.deleted_at.desc()).all():
                    items.append({"entity_type": entity_type, "id": item.id, "title": getattr(item, label_field),
                                  "deleted_at": _iso(item.deleted_at), "impact": self._trash_impact(entity_type, item)})
            return sorted(items, key=lambda value: value["deleted_at"] or "", reverse=True)
        return self._run(query)

    @staticmethod
    def _trash_model(entity_type):
        return {
            "project": ResearchProject, "lab_record": LabRecord, "literature": LiteratureItem,
            "library_item": LibraryItem, "note": Note, "task": Task, "weekly_report": WeeklyReport,
            "calendar_event": CalendarEvent,
        }.get(entity_type)

    @staticmethod
    def _trash_impact(entity_type, item):
        if entity_type == "project":
            return f"关联实验记录 {LabRecord.query.filter_by(project_id=item.id, is_deleted=False).count()} 条"
        if entity_type == "lab_record":
            return f"步骤 {len(item.steps)} 条，修订 {len(item.revisions)} 条"
        if entity_type == "library_item":
            return "托管文件永久删除时会清理本地副本；外部文件只清理索引"
        if entity_type == "calendar_event":
            return "仅移除手动日历事件；实验、任务和周报自动聚合项不受影响"
        return "永久删除后不可从回收站恢复"

    def trash_move(self, payload):
        entity_type = _string(payload.get("entity_type"), 30, "entity_type", required=True)
        entity_id = _positive_id(payload.get("id"))

        def move(workspace):
            model = self._trash_model(entity_type)
            if not model:
                raise ValidationError("对象类型无效。")
            item = model.query.filter_by(id=entity_id, workspace_id=workspace.id, is_deleted=False).first()
            if not item:
                raise NotFoundError("对象不存在或已在回收站。")
            item.is_deleted, item.deleted_at = True, utcnow()
            self._index_entity(workspace, {"lab_record": "record", "library_item": "file", "weekly_report": "weekly"}.get(entity_type, entity_type), item)
            return {"moved": True, "impact": self._trash_impact(entity_type, item)}
        return self._run(move)

    def trash_restore(self, payload):
        entity_type = _string(payload.get("entity_type"), 30, "entity_type", required=True)
        entity_id = _positive_id(payload.get("id"))

        def restore(workspace):
            model = self._trash_model(entity_type)
            if not model:
                raise ValidationError("对象类型无效。")
            item = model.query.filter_by(id=entity_id, workspace_id=workspace.id, is_deleted=True).first()
            if not item:
                raise NotFoundError("回收站中没有该对象。")
            item.is_deleted, item.deleted_at = False, None
            self._index_entity(workspace, {"lab_record": "record", "library_item": "file", "weekly_report": "weekly"}.get(entity_type, entity_type), item)
            return {"restored": True}
        return self._run(restore)

    def trash_purge(self, payload):
        entity_type = _string(payload.get("entity_type"), 30, "entity_type", required=True)
        entity_id = _positive_id(payload.get("id"))

        def purge(workspace):
            model = self._trash_model(entity_type)
            if not model:
                raise ValidationError("对象类型无效。")
            item = model.query.filter_by(id=entity_id, workspace_id=workspace.id, is_deleted=True).first()
            if not item:
                raise NotFoundError("回收站中没有该对象。")
            if entity_type == "library_item" and item.storage_mode == "managed":
                root = (Path(self.flask_app.instance_path) / "library").resolve()
                path = (root / str(item.managed_relative_path or "")).resolve()
                if root != path and root not in path.parents:
                    raise ValidationError("托管文件路径越界。")
                path.unlink(missing_ok=True)
            if entity_type == "library_item":
                FileOperation.query.filter_by(library_item_id=item.id).delete(synchronize_session=False)
            SearchDocument.query.filter_by(
                workspace_id=workspace.id,
                entity_type={"lab_record": "record", "library_item": "file", "weekly_report": "weekly"}.get(entity_type, entity_type),
                entity_id=item.id,
            ).delete(synchronize_session=False)
            db.session.delete(item)
            return {"purged": True}
        return self._run(purge)

    def settings_get(self, payload=None):
        def query(workspace):
            user_id = workspace.legacy_user_id
            setting = WorkspaceSetting.query.filter_by(user_id=user_id).first()
            preset = ApiPreset.query.filter_by(user_id=user_id, is_default=True).first()
            connection = ZoteroConnection.query.filter_by(workspace_id=workspace.id).first()
            return {
                "workspace": {"name": workspace.name, "timezone": workspace.timezone,
                              "data_path": str(Path(self.flask_app.instance_path).resolve())},
                "executors": [{"id": item.id, "name": item.name, "role": item.role, "is_active": item.is_active}
                              for item in Executor.query.filter_by(workspace_id=workspace.id).order_by(Executor.name)],
                "save": {"autosave": setting.execution_autosave if setting else True,
                         "interval": setting.execution_autosave_interval if setting else 30},
                "zotero": {"base_url": connection.base_url if connection else "http://127.0.0.1:23119",
                           "state": connection.connection_state if connection else "unknown"},
                "api": {"name": preset.name if preset else "未配置", "url": preset.api_url if preset else "",
                        "model": preset.text_model if preset else "", "enabled": bool(preset and preset.is_enabled),
                        "has_key": bool(preset and preset.encrypted_api_key)},
                "about": {"product": "R/LAB Research Assistant", "author": "面壁者"},
            }
        return self._run(query)

    def settings_save(self, payload):
        section = _string(payload.get("section"), 30, "section", required=True)

        def save(workspace):
            user_id = workspace.legacy_user_id
            if section == "workspace":
                workspace.name = _string(payload.get("name"), 120, "name", required=True)
                workspace.timezone = _string(payload.get("timezone") or "Asia/Shanghai", 64, "timezone")
            elif section == "save":
                setting = WorkspaceSetting.query.filter_by(user_id=user_id).first()
                if not setting:
                    setting = WorkspaceSetting(user_id=user_id)
                    db.session.add(setting)
                setting.execution_autosave = bool(payload.get("autosave"))
                interval = int(payload.get("interval") or 30)
                if interval < 5 or interval > 600:
                    raise ValidationError("自动保存间隔必须为 5 到 600 秒。")
                setting.execution_autosave_interval = interval
            elif section == "zotero":
                connection = ZoteroConnection.query.filter_by(workspace_id=workspace.id).first()
                if not connection:
                    connection = ZoteroConnection(workspace_id=workspace.id)
                    db.session.add(connection)
                base_url = _string(payload.get("base_url"), 500, "base_url", required=True).rstrip("/")
                if base_url != "http://127.0.0.1:23119":
                    raise ValidationError("Zotero Local API 地址必须为 http://127.0.0.1:23119。")
                connection.base_url = base_url
            elif section == "api":
                preset = ApiPreset.query.filter_by(user_id=user_id, is_default=True).first()
                if not preset:
                    preset = ApiPreset(user_id=user_id, name="桌面默认", is_default=True)
                    db.session.add(preset)
                preset.api_url = _string(payload.get("url"), 500, "url", required=True).rstrip("/")
                preset.text_model = _string(payload.get("model"), 160, "model", required=True)
                preset.is_enabled = bool(payload.get("enabled"))
                key = str(payload.get("api_key") or "").strip()
                if key:
                    preset.set_api_key(key)
            elif section == "executor":
                name = _string(payload.get("name"), 120, "name", required=True)
                normalized = " ".join(name.casefold().split())
                existing = Executor.query.filter_by(workspace_id=workspace.id, normalized_name=normalized).first()
                if existing:
                    existing.role, existing.is_active = _string(payload.get("role"), 120, "role"), True
                else:
                    db.session.add(Executor(user_id=user_id, workspace_id=workspace.id, name=name,
                                            normalized_name=normalized, role=_string(payload.get("role"), 120, "role")))
            else:
                raise ValidationError("设置分区无效。")
            return {"saved": True, "section": section}
        return self._run(save)

    def create_backup(self, payload=None):
        with self.flask_app.app_context():
            db.session.remove()
            source = Path(self.flask_app.instance_path) / "research.db"
            if not source.is_file():
                raise NotFoundError("当前数据库文件不存在。")
            target_dir = Path(self.flask_app.config["BACKUP_DIR"]).resolve()
            target_dir.mkdir(parents=True, exist_ok=True)
            target = target_dir / f"research-{datetime.now():%Y%m%d-%H%M%S}.db"
            with sqlite3.connect(source) as source_db, sqlite3.connect(target) as target_db:
                source_db.backup(target_db)
            return {"path": str(target), "size_bytes": target.stat().st_size}

    def export_record(self, payload):
        record_id = _positive_id(payload.get("id"))
        target = Path(_string(payload.get("path"), 2000, "path", required=True)).resolve()

        def export(workspace):
            record = LabRecord.query.filter_by(id=record_id, workspace_id=workspace.id, is_deleted=False).first()
            if not record:
                raise NotFoundError("实验记录不存在。")
            data = self._record_dto(record, detail=True)
            attachment_rows = db.session.execute(
                db.select(LibraryItem)
                .select_from(lab_record_library_item)
                .join(LibraryItem, LibraryItem.id == lab_record_library_item.c.library_item_id)
                .where(
                    lab_record_library_item.c.record_id == record.id,
                    LibraryItem.is_deleted.is_(False),
                )
                .order_by(LibraryItem.updated_at.desc())
            ).scalars().all()
            view = _record_export_view(record, attachment_rows)
            suffix = target.suffix.lower()
            target.parent.mkdir(parents=True, exist_ok=True)
            if suffix == ".json":
                target.write_text(json.dumps({**data, "export": view}, ensure_ascii=False, indent=2), encoding="utf-8")
            elif suffix in {".md", ".txt"}:
                _write_record_markdown_table_template(target, view)
            elif suffix == ".docx":
                _write_record_docx_table_template(target, view)
            elif suffix == ".pdf":
                font_path = self.flask_app.config.get("PDF_FONT_PATH")
                if not font_path:
                    raise ValidationError("未找到可用中文字体，PDF 导出不可用。")
                _write_record_pdf_table_template(target, view, font_path)
            else:
                raise ValidationError("仅支持导出 JSON、Markdown、Word 或 PDF。")
            return {"path": str(target), "size_bytes": target.stat().st_size}
        return self._run(export)

    def generate_weekly_ppt(self, payload):
        target = Path(_string(payload.get("path"), 2000, "path", required=True)).resolve()
        if target.suffix.lower() != ".pptx":
            raise ValidationError("PPT 生成位置必须使用 .pptx 扩展名。")
        record_ids = payload.get("record_ids") if isinstance(payload.get("record_ids"), list) else []
        record_ids = list(dict.fromkeys(_positive_id(value, "record_ids") for value in record_ids))[:100]
        if not record_ids:
            raise ValidationError("请至少选择一条实验记录。")

        def generate(workspace):
            records = LabRecord.query.filter(
                LabRecord.workspace_id == workspace.id, LabRecord.is_deleted.is_(False),
                LabRecord.id.in_(record_ids),
            ).order_by(LabRecord.experiment_date, LabRecord.id).all()
            if len(records) != len(record_ids):
                raise ValidationError("部分实验记录不存在或已移入回收站。")
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.util import Inches, Pt

            deck = Presentation()
            deck.slide_width = Inches(13.333)
            deck.slide_height = Inches(7.5)
            title_slide = deck.slides.add_slide(deck.slide_layouts[0])
            title_slide.shapes.title.text = _string(payload.get("title") or "研究进展与结果", 180, "title")
            title_slide.placeholders[1].text = f"R/LAB · {date.today():%Y-%m-%d}\n实验记录 {len(records)} 条"
            for record in records:
                slide = deck.slides.add_slide(deck.slide_layouts[5])
                slide.shapes.title.text = record.title
                box = slide.shapes.add_textbox(Inches(.75), Inches(1.35), Inches(11.8), Inches(5.3))
                frame = box.text_frame
                frame.clear()
                sections = (
                    ("研究目的", record.objective), ("实验设计", record.design),
                    ("实际结果", record.actual_result), ("分析与结论", f"{record.analysis}\n{record.conclusion}"),
                    ("下一步", record.next_steps),
                )
                for index, (heading, value) in enumerate(sections):
                    paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                    paragraph.text = f"{heading}  {value or '未填写'}"
                    paragraph.font.size = Pt(17 if index else 19)
                    paragraph.font.color.rgb = RGBColor(32, 33, 42)
                    paragraph.space_after = Pt(12)
            target.parent.mkdir(parents=True, exist_ok=True)
            deck.save(target)
            item = LibraryItem(
                workspace_id=workspace.id, display_name=target.name, original_name=target.name,
                description=f"由周报 PPT 工具生成，来源实验记录 {len(records)} 条。",
                kind="report", storage_mode="external", external_path=str(target),
                path_normalized=str(target).replace("/", "\\").casefold(),
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                size_bytes=target.stat().st_size, sha256=_sha256(target), link_status="available",
                ai_readability="metadata_only", last_verified_at=utcnow(),
            )
            db.session.add(item)
            db.session.flush()
            self._activity(workspace, "weekly_ppt_generated", f"生成 PPT：{target.name}", "library_item", item.id)
            return {"path": str(target), "size_bytes": target.stat().st_size, "library_item_id": item.id}
        return self._run(generate)
